from flask import Flask, render_template, request, jsonify
import io
from flask_cors import CORS
import requests
import json
import os
import re
from datetime import datetime
from bs4 import BeautifulSoup
import PyPDF2
from docx import Document
from pptx import Presentation
from google import genai
from prompt_loader import format_prompt

# Optional server-side defaults (env / config.py). Users normally enter their own in the UI.
CONFIG_AI_API_KEY = None
CONFIG_LLM_MODEL = None
try:
    import config as _cfg_mod
    _k = (getattr(_cfg_mod, "AI_API_KEY", None) or getattr(_cfg_mod, "GEMINI_API_KEY", None) or "")
    _k = str(_k).strip()
    if _k:
        CONFIG_AI_API_KEY = _k
        print("✓ Optional default AI API key loaded from config/env (users can override in the app)")
    _m = (getattr(_cfg_mod, "LLM_MODEL", None) or getattr(_cfg_mod, "AI_MODEL", None) or "")
    _m = str(_m).strip()
    if _m:
        CONFIG_LLM_MODEL = _m
        print("✓ Optional default LLM model loaded from config/env")
except ImportError:
    print("ℹ No config.py — add an AI API key in Configure Canvas, or set AI_API_KEY / GEMINI_API_KEY env var")
except Exception as e:
    print(f"⚠ Could not load optional AI API config: {e}")


def resolve_ai_api_key(request_json=None):
    """
    Resolve API key for LLM features (currently Google GenAI / Gemini-compatible).
    Priority: JSON ai_api_key / llm_api_key / gemini_api_key → env AI_API_KEY or GEMINI_API_KEY → config.
    """
    if request_json:
        for field in ("ai_api_key", "llm_api_key", "gemini_api_key"):
            user_key = (request_json.get(field) or "").strip()
            if user_key:
                return user_key
    for env_name in ("AI_API_KEY", "GEMINI_API_KEY"):
        env_key = os.environ.get(env_name, "").strip()
        if env_key:
            return env_key
    return CONFIG_AI_API_KEY


def resolve_llm_model(request_json=None):
    """
    Resolve LLM model id for generate_content.
    Priority: JSON llm_model / ai_model → env LLM_MODEL / AI_MODEL / GEMINI_MODEL → config.LLM_MODEL.
    """
    if request_json:
        for field in ("llm_model", "ai_model"):
            v = (request_json.get(field) or "").strip()
            if v:
                return v
    for env_name in ("LLM_MODEL", "AI_MODEL", "GEMINI_MODEL"):
        v = os.environ.get(env_name, "").strip()
        if v:
            return v
    return CONFIG_LLM_MODEL


app = Flask(__name__)
CORS(app)

# Canvas API base URL - USF Learn Canvas instance
CANVAS_BASE_URL = "https://usflearn.instructure.com"

# Data directory
DATA_DIR = "data/courses"
os.makedirs(DATA_DIR, exist_ok=True)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/material-selection')
def material_selection():
    return render_template('material_selection.html')

@app.route('/file-selection')
def file_selection():
    """File selection page for lecture materials"""
    return render_template('file_selection.html')

@app.route('/discussion-selection')
@app.route('/discussion-selection/')
def discussion_selection():
    """Discussion selection page for discussion prompts"""
    return render_template('discussion_selection.html')

@app.route('/all-materials-selection')
@app.route('/all-materials-selection/')
def all_materials_selection():
    """All materials selection page (excluding syllabus)."""
    return render_template('all_materials_selection.html')

@app.route('/comparison-report')
def comparison_report():
    return render_template('comparison_report.html')

@app.route('/api/compare-with-report', methods=['POST'])
def compare_with_report():
    """Compare courses and generate detailed Gemini report"""
    try:
        data = request.get_json()
        courses_a = data.get('courses_a', [])
        courses_b = data.get('courses_b', [])
        material_type = data.get('material_type', 'syllabus')
        
        if not courses_a or not courses_b:
            return jsonify({'error': 'Course IDs are required'}), 400
        
        print(f"\n=== Generating Detailed Comparison Report ===")
        print(f"Course A: {courses_a[0]}, Course B: {courses_b[0]}")
        print(f"Material Type: {material_type}")
        
        # Load course data
        course_a_data = load_course_data(courses_a[0])
        course_b_data = load_course_data(courses_b[0])
        
        if not course_a_data or not course_b_data:
            return jsonify({'error': 'Failed to load course data'}), 404

        ai_key = resolve_ai_api_key(data)
        if not ai_key:
            return jsonify({
                'error': 'AI API key is required. Add your provider API key under Configure Canvas, or set AI_API_KEY (or GEMINI_API_KEY) on the server.'
            }), 400

        llm_model = resolve_llm_model(data)
        if not llm_model:
            return jsonify({
                'error': 'LLM model name is required. Enter it under Configure Canvas, or set LLM_MODEL (or AI_MODEL / GEMINI_MODEL) on the server.'
            }), 400
        
        # Generate detailed report with Gemini
        report = generate_detailed_comparison_report(
            course_a_data,
            course_b_data,
            material_type,
            ai_key,
            llm_model
        )
        
        return jsonify({
            'success': True,
            'course_a': {
                'id': courses_a[0],
                'name': course_a_data.get('metadata', {}).get('course_name', 'Course A'),
                'term': course_a_data.get('metadata', {}).get('course_term', '')
            },
            'course_b': {
                'id': courses_b[0],
                'name': course_b_data.get('metadata', {}).get('course_name', 'Course B'),
                'term': course_b_data.get('metadata', {}).get('course_term', '')
            },
            'material_type': material_type,
            'report': report
        }), 200
        
    except Exception as e:
        print(f"Error generating report: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

def generate_detailed_comparison_report(course_a_data, course_b_data, material_type, api_key, llm_model):
    """Use configured LLM to generate a detailed comparison report"""
    try:
        from google import genai
        
        client = genai.Client(api_key=api_key)
        
        # Get course names
        course_a_name = course_a_data.get('metadata', {}).get('course_name', 'Course A')
        course_b_name = course_b_data.get('metadata', {}).get('course_name', 'Course B')
        
        # For SYLLABUS comparison, use the FILTERED JSON data
        if material_type == 'syllabus':
            syllabi_a = course_a_data.get('syllabus', [])
            syllabi_b = course_b_data.get('syllabus', [])
            
            if not syllabi_a or not syllabi_b:
                raise Exception("Syllabus data not found for one or both courses")
            
            # Get the filtered data (description, outcomes, topics/schedule)
            filtered_a = syllabi_a[0].get('filtered_data', {})
            filtered_b = syllabi_b[0].get('filtered_data', {})
            
            # Prepare structured JSON for comparison (ONLY extracted content, no additions)
            course_a_json = {
                "course_name": course_a_name,
                "course_description": filtered_a.get('course_description', ''),
                "learning_outcomes": filtered_a.get('learning_outcomes', ''),
                "topics_schedule": filtered_a.get('topics_schedule', '')
            }
            
            course_b_json = {
                "course_name": course_b_name,
                "course_description": filtered_b.get('course_description', ''),
                "learning_outcomes": filtered_b.get('learning_outcomes', ''),
                "topics_schedule": filtered_b.get('topics_schedule', '')
            }
            
            # Convert to formatted JSON strings
            content_a = json.dumps(course_a_json, indent=2)
            content_b = json.dumps(course_b_json, indent=2)
        
        # For LECTURES comparison, use the COMBINED JSON (all lecture files in one)
        elif material_type == 'lectures':
            lectures_a = course_a_data.get('lectures_combined')
            lectures_b = course_b_data.get('lectures_combined')
            
            if not lectures_a or not lectures_b:
                raise Exception("Combined lecture data not found for one or both courses. Please re-extract courses.")
            
            # Use the combined content directly
            content_a = f"Course: {course_a_name}\n"
            content_a += f"Total Lecture Files: {lectures_a.get('total_files', 0)}\n"
            content_a += f"Files: {', '.join([f['name'] for f in lectures_a.get('file_index', [])])}\n\n"
            content_a += lectures_a.get('combined_content', '')
            
            content_b = f"Course: {course_b_name}\n"
            content_b += f"Total Lecture Files: {lectures_b.get('total_files', 0)}\n"
            content_b += f"Files: {', '.join([f['name'] for f in lectures_b.get('file_index', [])])}\n\n"
            content_b += lectures_b.get('combined_content', '')
        
        # For GRADED ASSIGNMENTS (tests), use the combined JSON (selected graded assignments in one)
        elif material_type == 'graded_assignments':
            ga_a = course_a_data.get('graded_assignments_combined')
            ga_b = course_b_data.get('graded_assignments_combined')
            if not ga_a or not ga_b:
                raise Exception("Graded assignments combined data not found for one or both courses. Please re-extract.")
            content_a = f"Course: {course_a_name}\n"
            content_a += f"Total Assignments: {ga_a.get('total_assignments', 0)}\n"
            content_a += f"Assignments: {', '.join([a['name'] for a in ga_a.get('assignment_index', [])])}\n\n"
            content_a += ga_a.get('combined_content', '')
            content_b = f"Course: {course_b_name}\n"
            content_b += f"Total Assignments: {ga_b.get('total_assignments', 0)}\n"
            content_b += f"Assignments: {', '.join([a['name'] for a in ga_b.get('assignment_index', [])])}\n\n"
            content_b += ga_b.get('combined_content', '')
        
        elif material_type == 'discussions':
            d_a = course_a_data.get('discussions_combined')
            d_b = course_b_data.get('discussions_combined')
            if not d_a or not d_b:
                raise Exception("Discussion combined data not found for one or both courses. Please re-extract.")
            content_a = f"Course: {course_a_name}\n"
            content_a += f"Total Discussion Posts: {d_a.get('total_discussions', 0)}\n"
            content_a += f"Posts: {', '.join([a['name'] for a in d_a.get('discussion_index', [])])}\n\n"
            content_a += d_a.get('combined_content', '')
            content_b = f"Course: {course_b_name}\n"
            content_b += f"Total Discussion Posts: {d_b.get('total_discussions', 0)}\n"
            content_b += f"Posts: {', '.join([a['name'] for a in d_b.get('discussion_index', [])])}\n\n"
            content_b += d_b.get('combined_content', '')
        
        elif material_type == 'all_selected':
            all_a = course_a_data.get('all_selected_combined')
            all_b = course_b_data.get('all_selected_combined')
            if not all_a or not all_b:
                raise Exception("All selected materials combined data not found for one or both courses. Please re-extract.")
            content_a = f"Course: {course_a_name}\n"
            content_a += f"Total Selected Items: {all_a.get('total_items', 0)}\n\n"
            content_a += all_a.get('combined_content', '')
            content_b = f"Course: {course_b_name}\n"
            content_b += f"Total Selected Items: {all_b.get('total_items', 0)}\n\n"
            content_b += all_b.get('combined_content', '')
            
        else:
            # For other material types, use regular content
            if material_type == 'assignments':
                materials_a = course_a_data.get('assignments', [])
                materials_b = course_b_data.get('assignments', [])
            elif material_type == 'quizzes':
                materials_a = course_a_data.get('quizzes', [])
                materials_b = course_b_data.get('quizzes', [])
            elif material_type == 'discussions':
                materials_a = course_a_data.get('discussions', [])
                materials_b = course_b_data.get('discussions', [])
            else:
                materials_a = []
                materials_b = []
            
            content_a = '\n\n'.join([f"{m.get('name', 'Untitled')}:\n{m.get('content', '')}" for m in materials_a])
            content_b = '\n\n'.join([f"{m.get('name', 'Untitled')}:\n{m.get('content', '')}" for m in materials_b])
        
        # Create specific prompt based on material type
        if material_type == 'syllabus':
            prompt = format_prompt('syllabus_comparison',
                course_a_name=course_a_name,
                course_b_name=course_b_name,
                content_a=content_a,
                content_b=content_b
            )
        
        elif material_type == 'lectures':
            prompt = format_prompt('lectures_comparison',
                course_a_name=course_a_name,
                course_b_name=course_b_name,
                content_a=content_a,
                content_b=content_b
            )
        
        elif material_type == 'graded_assignments':
            prompt = format_prompt('generic_comparison',
                course_a_name=course_a_name,
                course_b_name=course_b_name,
                content_a=content_a,
                content_b=content_b
            )
        
        elif material_type == 'discussions':
            prompt = format_prompt('generic_comparison',
                course_a_name=course_a_name,
                course_b_name=course_b_name,
                content_a=content_a,
                content_b=content_b
            )
        
        elif material_type == 'all_selected':
            prompt = format_prompt('generic_comparison',
                course_a_name=course_a_name,
                course_b_name=course_b_name,
                content_a=content_a,
                content_b=content_b
            )
        
        else:
            # Generic prompt for other material types
            prompt = format_prompt('generic_comparison',
                course_a_name=course_a_name,
                course_b_name=course_b_name,
                content_a=content_a,
                content_b=content_b
            )

        print(f"Calling LLM ({llm_model}) to generate detailed report...")
        
        response = client.models.generate_content(
            model=llm_model,
            contents=prompt
        )
        
        report_text = response.text
        
        print("✓ Detailed report generated successfully")
        
        # Calculate material counts
        if material_type == 'syllabus':
            material_count_a = 1 if course_a_data.get('syllabus') else 0
            material_count_b = 1 if course_b_data.get('syllabus') else 0
        elif material_type == 'graded_assignments':
            ga_a = course_a_data.get('graded_assignments_combined') or {}
            ga_b = course_b_data.get('graded_assignments_combined') or {}
            material_count_a = ga_a.get('total_assignments', 0)
            material_count_b = ga_b.get('total_assignments', 0)
        elif material_type == 'discussions':
            d_a = course_a_data.get('discussions_combined') or {}
            d_b = course_b_data.get('discussions_combined') or {}
            material_count_a = d_a.get('total_discussions', 0)
            material_count_b = d_b.get('total_discussions', 0)
        elif material_type == 'all_selected':
            all_a = course_a_data.get('all_selected_combined') or {}
            all_b = course_b_data.get('all_selected_combined') or {}
            material_count_a = all_a.get('total_items', 0)
            material_count_b = all_b.get('total_items', 0)
        else:
            materials_a_list = materials_a if 'materials_a' in locals() else []
            materials_b_list = materials_b if 'materials_b' in locals() else []
            material_count_a = len(materials_a_list)
            material_count_b = len(materials_b_list)
        
        return {
            'report_text': report_text,
            'generated_at': datetime.now().isoformat(),
            'materials_compared': {
                'course_a_count': material_count_a,
                'course_b_count': material_count_b
            }
        }
        
    except Exception as e:
        print(f"Error generating detailed report: {e}")
        import traceback
        traceback.print_exc()
        return {
            'report_text': f"Error generating report: {str(e)}",
            'error': str(e),
            'generated_at': datetime.now().isoformat()
        }

@app.route('/extraction')
def extraction():
    return render_template('extraction.html')

@app.route('/api/compare-courses', methods=['POST'])
def compare_courses():
    """Compare materials between two courses using Gemini LLM"""
    try:
        data = request.get_json()
        courses_a = data.get('courses_a', [])
        courses_b = data.get('courses_b', [])
        
        if not courses_a or not courses_b:
            return jsonify({'error': 'Both course sets are required'}), 400

        ai_key = resolve_ai_api_key(data)
        if not ai_key:
            return jsonify({
                'error': 'AI API key is required. Add your key in Configure Canvas or set AI_API_KEY / GEMINI_API_KEY on the server.'
            }), 400

        llm_model = resolve_llm_model(data)
        if not llm_model:
            return jsonify({
                'error': 'LLM model name is required. Enter it under Configure Canvas or set LLM_MODEL on the server.'
            }), 400
        
        print(f"\n=== Starting course comparison ===")
        print(f"Course Set A: {courses_a}")
        print(f"Course Set B: {courses_b}")
        
        # Load extracted course data
        course_a_data = load_course_data(courses_a[0])  # For now, compare first course from each set
        course_b_data = load_course_data(courses_b[0])
        
        if not course_a_data or not course_b_data:
            return jsonify({'error': 'Could not load course data. Please extract courses first.'}), 404
        
        # Import comparison module
        from compare import compare_course_materials, calculate_overall_similarity
        
        # Create progress directory
        progress_dir = os.path.join(DATA_DIR, 'comparison_progress')
        os.makedirs(progress_dir, exist_ok=True)
        
        # Compare materials by type with progress saving
        # IMPORTANT: Separate graded vs non-graded materials
        comparison_results = {}
        
        # Helper function to separate graded vs non-graded
        def separate_by_graded(materials):
            graded = [m for m in materials if m.get('graded', False)]
            non_graded = [m for m in materials if not m.get('graded', False)]
            return graded, non_graded
        
        # Compare assignments (GRADED only)
        if course_a_data.get('assignments') and course_b_data.get('assignments'):
            graded_a, non_graded_a = separate_by_graded(course_a_data['assignments'])
            graded_b, non_graded_b = separate_by_graded(course_b_data['assignments'])
            
            # Compare GRADED assignments
            if graded_a and graded_b:
                print(f"\nComparing GRADED assignments...")
                print(f"  Course A: {len(graded_a)} graded assignments")
                print(f"  Course B: {len(graded_b)} graded assignments")
                progress_file = os.path.join(progress_dir, f'assignments_graded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['assignments_graded'] = compare_course_materials(
                    graded_a,
                    graded_b,
                    'graded_assignment',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
            
            # Compare NON-GRADED assignments
            if non_graded_a and non_graded_b:
                print(f"\nComparing NON-GRADED assignments...")
                print(f"  Course A: {len(non_graded_a)} non-graded assignments")
                print(f"  Course B: {len(non_graded_b)} non-graded assignments")
                progress_file = os.path.join(progress_dir, f'assignments_nongraded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['assignments_nongraded'] = compare_course_materials(
                    non_graded_a,
                    non_graded_b,
                    'non_graded_assignment',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
        
        # Compare quizzes (GRADED only - quizzes are typically graded)
        if course_a_data.get('quizzes') and course_b_data.get('quizzes'):
            graded_a, non_graded_a = separate_by_graded(course_a_data['quizzes'])
            graded_b, non_graded_b = separate_by_graded(course_b_data['quizzes'])
            
            # Compare GRADED quizzes
            if graded_a and graded_b:
                print(f"\nComparing GRADED quizzes...")
                print(f"  Course A: {len(graded_a)} graded quizzes")
                print(f"  Course B: {len(graded_b)} graded quizzes")
                progress_file = os.path.join(progress_dir, f'quizzes_graded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['quizzes_graded'] = compare_course_materials(
                    graded_a,
                    graded_b,
                    'graded_quiz',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
            
            # Compare NON-GRADED quizzes (practice quizzes, surveys, etc.)
            if non_graded_a and non_graded_b:
                print(f"\nComparing NON-GRADED quizzes...")
                print(f"  Course A: {len(non_graded_a)} non-graded quizzes")
                print(f"  Course B: {len(non_graded_b)} non-graded quizzes")
                progress_file = os.path.join(progress_dir, f'quizzes_nongraded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['quizzes_nongraded'] = compare_course_materials(
                    non_graded_a,
                    non_graded_b,
                    'non_graded_quiz',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
        
        # Compare files (NON-GRADED - lecture slides, readings, etc.)
        if course_a_data.get('files') and course_b_data.get('files'):
            graded_a, non_graded_a = separate_by_graded(course_a_data['files'])
            graded_b, non_graded_b = separate_by_graded(course_b_data['files'])
            
            # Files are typically non-graded (lecture slides, readings)
            if non_graded_a and non_graded_b:
                print(f"\nComparing NON-GRADED files (lecture slides, readings)...")
                print(f"  Course A: {len(non_graded_a)} files")
                print(f"  Course B: {len(non_graded_b)} files")
                progress_file = os.path.join(progress_dir, f'files_nongraded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['files_nongraded'] = compare_course_materials(
                    non_graded_a,
                    non_graded_b,
                    'lecture_material',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
            
            # If any files are graded (rare, but possible)
            if graded_a and graded_b:
                print(f"\nComparing GRADED files...")
                print(f"  Course A: {len(graded_a)} graded files")
                print(f"  Course B: {len(graded_b)} graded files")
                progress_file = os.path.join(progress_dir, f'files_graded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['files_graded'] = compare_course_materials(
                    graded_a,
                    graded_b,
                    'graded_file',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
        
        # Compare SYLLABUS (syllabus-to-syllabus only)
        if course_a_data.get('syllabus') and course_b_data.get('syllabus'):
            if len(course_a_data['syllabus']) > 0 and len(course_b_data['syllabus']) > 0:
                print(f"\nComparing SYLLABUS...")
                print(f"  Course A: {len(course_a_data['syllabus'])} syllabus")
                print(f"  Course B: {len(course_b_data['syllabus'])} syllabus")
                progress_file = os.path.join(progress_dir, f'syllabus_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['syllabus'] = compare_course_materials(
                    course_a_data['syllabus'],
                    course_b_data['syllabus'],
                    'syllabus',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
        
        # Compare INDIVIDUAL QUIZ QUESTIONS (question-to-question)
        if course_a_data.get('quiz_questions') and course_b_data.get('quiz_questions'):
            graded_a, non_graded_a = separate_by_graded(course_a_data['quiz_questions'])
            graded_b, non_graded_b = separate_by_graded(course_b_data['quiz_questions'])
            
            # Compare GRADED quiz questions
            if graded_a and graded_b:
                print(f"\nComparing GRADED quiz questions (question-to-question)...")
                print(f"  Course A: {len(graded_a)} quiz questions")
                print(f"  Course B: {len(graded_b)} quiz questions")
                progress_file = os.path.join(progress_dir, f'quiz_questions_graded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['quiz_questions_graded'] = compare_course_materials(
                    graded_a,
                    graded_b,
                    'graded_quiz_question',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
            
            # Compare NON-GRADED quiz questions (practice questions, survey questions)
            if non_graded_a and non_graded_b:
                print(f"\nComparing NON-GRADED quiz questions...")
                print(f"  Course A: {len(non_graded_a)} quiz questions")
                print(f"  Course B: {len(non_graded_b)} quiz questions")
                progress_file = os.path.join(progress_dir, f'quiz_questions_nongraded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['quiz_questions_nongraded'] = compare_course_materials(
                    non_graded_a,
                    non_graded_b,
                    'non_graded_quiz_question',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
        
        # Compare discussions (can be graded or non-graded)
        if course_a_data.get('discussions') and course_b_data.get('discussions'):
            graded_a, non_graded_a = separate_by_graded(course_a_data['discussions'])
            graded_b, non_graded_b = separate_by_graded(course_b_data['discussions'])
            
            # Compare GRADED discussions
            if graded_a and graded_b:
                print(f"\nComparing GRADED discussions...")
                print(f"  Course A: {len(graded_a)} graded discussions")
                print(f"  Course B: {len(graded_b)} graded discussions")
                progress_file = os.path.join(progress_dir, f'discussions_graded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['discussions_graded'] = compare_course_materials(
                    graded_a,
                    graded_b,
                    'graded_discussion',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
            
            # Compare NON-GRADED discussions
            if non_graded_a and non_graded_b:
                print(f"\nComparing NON-GRADED discussions...")
                print(f"  Course A: {len(non_graded_a)} non-graded discussions")
                print(f"  Course B: {len(non_graded_b)} non-graded discussions")
                progress_file = os.path.join(progress_dir, f'discussions_nongraded_{courses_a[0]}_{courses_b[0]}.json')
                comparison_results['discussions_nongraded'] = compare_course_materials(
                    non_graded_a,
                    non_graded_b,
                    'non_graded_discussion',
                    ai_key,
                    batch_size=10,
                    progress_file=progress_file,
                    parallel_workers=5,
                    llm_model=llm_model
                )
        
        # Calculate overall similarity
        overall_results = calculate_overall_similarity(comparison_results)
        
        print(f"\n=== Comparison complete ===")
        print(f"Overall similarity: {overall_results.get('overall_similarity', 0):.1f}%")
        
        return jsonify({
            'success': True,
            'course_a_name': course_a_data.get('metadata', {}).get('course_name', 'Course A'),
            'course_b_name': course_b_data.get('metadata', {}).get('course_name', 'Course B'),
            'comparison_results': comparison_results,
            'overall_results': overall_results
        }), 200
        
    except Exception as e:
        print(f"Error in comparison: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

def load_course_data(course_id):
    """Load all extracted data for a course"""
    try:
        course_dir = os.path.join(DATA_DIR, f"course_{course_id}")
        
        if not os.path.exists(course_dir):
            print(f"Course directory not found: {course_dir}")
            return None
        
        # Load metadata
        metadata_path = os.path.join(course_dir, "metadata.json")
        metadata = {}
        if os.path.exists(metadata_path):
            with open(metadata_path, 'r', encoding='utf-8') as f:
                metadata = json.load(f)
        
        # Load all materials
        # Load syllabus from syllabus/ folder (filtered + raw) or legacy root
        syllabus_dir = os.path.join(course_dir, "syllabus")
        syllabus_filtered_file = os.path.join(syllabus_dir, "syllabus_filtered.json")
        if not os.path.exists(syllabus_filtered_file):
            syllabus_filtered_file = os.path.join(course_dir, "syllabus_filtered.json")
        syllabus_data = []
        if os.path.exists(syllabus_filtered_file):
            try:
                with open(syllabus_filtered_file, 'r', encoding='utf-8') as f:
                    syllabus_data = [json.load(f)]
                print(f"  Loaded filtered syllabus")
            except Exception as e:
                print(f"Error loading filtered syllabus: {e}")
        elif os.path.exists(os.path.join(course_dir, "syllabus.json")):
            try:
                with open(os.path.join(course_dir, "syllabus.json"), 'r', encoding='utf-8') as f:
                    syllabus_data = [json.load(f)]
                print(f"  Loaded syllabus (old format)")
            except Exception as e:
                print(f"Error loading syllabus: {e}")
        
        # Load lectures from lecture_materials/ folder or legacy root
        lecture_materials_dir = os.path.join(course_dir, "lecture_materials")
        lectures_combined_file = os.path.join(lecture_materials_dir, "lectures_combined.json")
        if not os.path.exists(lectures_combined_file):
            lectures_combined_file = os.path.join(course_dir, "lectures_combined.json")
        lectures_combined = None
        if os.path.exists(lectures_combined_file):
            try:
                with open(lectures_combined_file, 'r', encoding='utf-8') as f:
                    lectures_combined = json.load(f)
                print(f"  Loaded combined lectures ({lectures_combined.get('total_files', 0)} files)")
            except Exception as e:
                print(f"Error loading combined lectures: {e}")
        
        # Load graded assignments from graded_assignments/ folder or legacy root
        graded_dir = os.path.join(course_dir, "graded_assignments")
        graded_assignments_combined_file = os.path.join(graded_dir, "graded_assignments_combined.json")
        if not os.path.exists(graded_assignments_combined_file):
            graded_assignments_combined_file = os.path.join(course_dir, "graded_assignments_combined.json")
        graded_assignments_combined = None
        if os.path.exists(graded_assignments_combined_file):
            try:
                with open(graded_assignments_combined_file, 'r', encoding='utf-8') as f:
                    graded_assignments_combined = json.load(f)
                print(f"  Loaded graded assignments combined ({graded_assignments_combined.get('total_assignments', 0)} assignments)")
            except Exception as e:
                print(f"Error loading graded assignments combined: {e}")

        # Load discussions from discussions/ folder
        discussions_dir = os.path.join(course_dir, "discussions")
        discussions_combined_file = os.path.join(discussions_dir, "discussions_combined.json")
        if not os.path.exists(discussions_combined_file):
            discussions_combined_file = os.path.join(course_dir, "discussions_combined.json")
        discussions_combined = None
        if os.path.exists(discussions_combined_file):
            try:
                with open(discussions_combined_file, 'r', encoding='utf-8') as f:
                    discussions_combined = json.load(f)
                print(f"  Loaded discussions combined ({discussions_combined.get('total_discussions', 0)} posts)")
            except Exception as e:
                print(f"Error loading discussions combined: {e}")

        # Load all-selected combined from all_materials/ folder
        all_materials_dir = os.path.join(course_dir, "all_materials")
        all_selected_combined_file = os.path.join(all_materials_dir, "all_selected_combined.json")
        if not os.path.exists(all_selected_combined_file):
            all_selected_combined_file = os.path.join(course_dir, "all_selected_combined.json")
        all_selected_combined = None
        if os.path.exists(all_selected_combined_file):
            try:
                with open(all_selected_combined_file, 'r', encoding='utf-8') as f:
                    all_selected_combined = json.load(f)
                print(f"  Loaded all materials combined ({all_selected_combined.get('total_items', 0)} items)")
            except Exception as e:
                print(f"Error loading all materials combined: {e}")
        
        data = {
            'metadata': metadata,
            'syllabus': syllabus_data,
            'lectures_combined': lectures_combined,
            'graded_assignments_combined': graded_assignments_combined,
            'discussions_combined': discussions_combined,
            'all_selected_combined': all_selected_combined,
            'assignments': load_materials_from_folder(os.path.join(course_dir, "assignments")),
            'quizzes': load_materials_from_folder(os.path.join(course_dir, "quizzes")),
            'quiz_questions': load_materials_from_folder(os.path.join(course_dir, "quiz_questions")),
            'files': load_materials_from_folder(os.path.join(course_dir, "files")),
            'discussions': load_materials_from_folder(os.path.join(course_dir, "discussions"))
        }
        # Avoid treating the combined discussions JSON as an individual discussion item.
        data['discussions'] = [
            d for d in data.get('discussions', [])
            if d.get('type') != 'discussions_combined'
        ]
        
        return data
        
    except Exception as e:
        print(f"Error loading course data: {e}")
        return None

def load_materials_from_folder(folder_path):
    """Load all JSON files from a folder"""
    materials = []
    
    if not os.path.exists(folder_path):
        return materials
    
    try:
        for filename in os.listdir(folder_path):
            if filename.endswith('.json'):
                file_path = os.path.join(folder_path, filename)
                with open(file_path, 'r', encoding='utf-8') as f:
                    material = json.load(f)
                    materials.append(material)
        
        return materials
        
    except Exception as e:
        print(f"Error loading materials from {folder_path}: {e}")
        return materials

@app.route('/api/files-list', methods=['POST'])
def get_files_list():
    """Fetch list of files from Canvas (metadata only, no content extraction)"""
    try:
        data = request.get_json()
        course_id = data.get('course_id')
        api_token = data.get('api_token')
        canvas_url = data.get('canvas_url', CANVAS_BASE_URL)
        
        if not course_id or not api_token:
            return jsonify({'error': 'Course ID and API token are required'}), 400
        
        headers = {
            'Authorization': f'Bearer {api_token}'
        }
        
        url = f"{canvas_url}/api/v1/courses/{course_id}/files"
        params = {'per_page': 100}
        response = requests.get(url, headers=headers, params=params, timeout=10)
        
        if response.status_code == 200:
            files = response.json()
            # File IDs linked from graded assignments (exclude these — lecture list = non-graded only)
            graded_file_ids = _graded_file_ids_for_course(course_id, canvas_url, headers)
            
            # Filter for document types only (PDF, PPT, Word) and exclude graded-assignment files
            document_files = []
            for file_info in files:
                file_id = file_info.get('id')
                if file_id in graded_file_ids:
                    continue
                file_name = file_info.get('display_name', '')
                mime_type = file_info.get('content-type', '')
                file_size = file_info.get('size', 0)
                
                # Check if it's a document type
                is_document = (
                    'pdf' in (mime_type or '').lower() or file_name.lower().endswith('.pdf') or
                    file_name.lower().endswith('.pptx') or
                    file_name.lower().endswith('.docx')
                )
                
                if is_document and file_size <= 50 * 1024 * 1024:  # Skip files > 50MB
                    document_files.append({
                        'id': file_id,
                        'name': file_name,
                        'size': file_size,
                        'size_mb': round(file_size / (1024 * 1024), 2),
                        'type': 'pdf' if '.pdf' in file_name.lower() else ('pptx' if '.pptx' in file_name.lower() else 'docx'),
                        'url': file_info.get('url')
                    })
            
            return jsonify({
                'success': True,
                'files': document_files,
                'total': len(document_files)
            }), 200
        else:
            return jsonify({'error': f'Failed to fetch files: {response.status_code}'}), response.status_code
            
    except Exception as e:
        print(f"Error fetching files list: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/assignments-list', methods=['POST'])
def get_assignments_list():
    """Fetch list of graded assignments from Canvas (metadata only). Uses is_graded: grading_type != 'not_graded' and points_possible > 0."""
    try:
        data = request.get_json()
        course_id = data.get('course_id')
        api_token = data.get('api_token')
        canvas_url = data.get('canvas_url', CANVAS_BASE_URL)
        
        if not course_id or not api_token:
            return jsonify({'error': 'Course ID and API token are required'}), 400
        
        headers = {'Authorization': f'Bearer {api_token}'}
        url = f"{canvas_url}/api/v1/courses/{course_id}/assignments"
        params = {'per_page': 100}
        response = requests.get(url, headers=headers, params=params, timeout=10)
        
        if response.status_code != 200:
            return jsonify({'error': f'Failed to fetch assignments: {response.status_code}'}), response.status_code
        
        assignments = response.json()
        skip_keywords = ['survey', 'participation', 'attendance', 'pitch', 'introductions', 'introduction', 'flyer']
        graded = []
        for a in assignments:
            if not is_graded(a):
                continue
            name_lower = (a.get('name') or '').lower()
            if any(kw in name_lower for kw in skip_keywords):
                continue
            graded.append({
                'id': a.get('id'),
                'name': a.get('name', 'Untitled'),
                'points_possible': a.get('points_possible', 0),
                'grading_type': a.get('grading_type', ''),
                'due_at': a.get('due_at'),
                'position': a.get('position', 999)
            })
        graded.sort(key=lambda x: (x['position'], x.get('due_at') or ''))
        return jsonify({'success': True, 'assignments': graded, 'total': len(graded)}), 200
    except Exception as e:
        print(f"Error fetching assignments list: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/discussions-list', methods=['POST'])
def get_discussions_list():
    """Fetch professor-created discussion posts (topic prompts only, no replies)."""
    try:
        data = request.get_json()
        course_id = data.get('course_id')
        api_token = data.get('api_token')
        canvas_url = data.get('canvas_url', CANVAS_BASE_URL)

        if not course_id or not api_token:
            return jsonify({'error': 'Course ID and API token are required'}), 400

        headers = {'Authorization': f'Bearer {api_token}'}
        discussions = list_professor_discussion_topics(course_id, canvas_url, headers)

        result = []
        for d in discussions:
            title = d.get('title', 'Untitled Discussion')
            result.append({
                'id': d.get('id'),
                'title': title,
                'author_name': (d.get('author') or {}).get('display_name', ''),
                'posted_at': d.get('posted_at') or d.get('created_at'),
                'points_possible': (d.get('assignment') or {}).get('points_possible', 0),
                'graded': d.get('assignment') is not None
            })

        result.sort(key=lambda x: (x.get('posted_at') or '', x.get('title') or ''))
        return jsonify({'success': True, 'discussions': result, 'total': len(result)}), 200
    except Exception as e:
        print(f"Error fetching discussions list: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/assignment-selection')
@app.route('/assignment-selection/')
def assignment_selection():
    """Assignment selection page for graded materials"""
    return render_template('assignment_selection.html')


@app.route('/api/prompts/default', methods=['GET'])
def get_default_prompts():
    """Get all default prompts"""
    try:
        prompts = {}
        prompts_dir = os.path.join(os.path.dirname(__file__), 'prompts')
        
        for filename in os.listdir(prompts_dir):
            if filename.endswith('.json'):
                prompt_name = filename[:-5]  # Remove .json extension
                prompt_path = os.path.join(prompts_dir, filename)
                with open(prompt_path, 'r', encoding='utf-8') as f:
                    prompts[prompt_name] = json.load(f)
        
        return jsonify(prompts), 200
    except Exception as e:
        print(f"Error loading default prompts: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/prompts/custom', methods=['GET'])
def get_custom_prompts():
    """Get all custom prompts"""
    try:
        custom_prompts_path = os.path.join(os.path.dirname(__file__), 'custom_prompts.json')
        if os.path.exists(custom_prompts_path):
            with open(custom_prompts_path, 'r', encoding='utf-8') as f:
                return jsonify(json.load(f)), 200
        return jsonify({}), 200
    except Exception as e:
        print(f"Error loading custom prompts: {e}")
        return jsonify({'error': str(e)}), 500

PROMPT_HISTORY_PATH = os.path.join(os.path.dirname(__file__), 'prompt_history.json')
PROMPT_HISTORY_MAX_VERSIONS = 20

def load_prompt_history():
    """Load prompt version history from file"""
    if os.path.exists(PROMPT_HISTORY_PATH):
        try:
            with open(PROMPT_HISTORY_PATH, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"Warning: Could not load prompt history: {e}")
    return {}

def save_prompt_history(history):
    """Write prompt version history to file"""
    with open(PROMPT_HISTORY_PATH, 'w', encoding='utf-8') as f:
        json.dump(history, f, indent=2, ensure_ascii=False)

def append_to_history(prompt_name, template, custom_prompts):
    """Append current version of a prompt to history before overwriting"""
    history = load_prompt_history()
    versions = history.get(prompt_name, [])

    # Determine the next version number
    next_version = (versions[0]['version'] + 1) if versions else 1

    # If there's already a saved custom version, push it into history
    if prompt_name in custom_prompts:
        existing = custom_prompts[prompt_name]
        versions.insert(0, {
            'version': next_version,
            'template': existing['template'],
            'saved_at': existing['saved_at']
        })
        # Cap history length
        versions = versions[:PROMPT_HISTORY_MAX_VERSIONS]

    history[prompt_name] = versions
    save_prompt_history(history)

@app.route('/api/prompts/save', methods=['POST'])
def save_custom_prompt():
    """Save a custom prompt (stored in custom_prompts.json), with version history"""
    try:
        data = request.get_json()
        prompt_name = data.get('prompt_name')
        template = data.get('template')
        
        if not prompt_name or not template:
            return jsonify({'error': 'prompt_name and template are required'}), 400
        
        # Load existing custom prompts
        custom_prompts_path = os.path.join(os.path.dirname(__file__), 'custom_prompts.json')
        custom_prompts = {}
        if os.path.exists(custom_prompts_path):
            with open(custom_prompts_path, 'r', encoding='utf-8') as f:
                custom_prompts = json.load(f)
        
        # Archive the current version before overwriting
        append_to_history(prompt_name, template, custom_prompts)

        # Save new custom prompt
        custom_prompts[prompt_name] = {
            'template': template,
            'saved_at': datetime.now().isoformat()
        }
        
        with open(custom_prompts_path, 'w', encoding='utf-8') as f:
            json.dump(custom_prompts, f, indent=2, ensure_ascii=False)
        
        return jsonify({'success': True, 'message': f'Prompt "{prompt_name}" saved'}), 200
    except Exception as e:
        print(f"Error saving custom prompt: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/prompts/history', methods=['GET'])
def get_prompt_history():
    """Return version history for a specific prompt"""
    try:
        prompt_name = request.args.get('prompt_name')
        if not prompt_name:
            return jsonify({'error': 'prompt_name query param is required'}), 400
        history = load_prompt_history()
        versions = history.get(prompt_name, [])
        return jsonify(versions), 200
    except Exception as e:
        print(f"Error fetching prompt history: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/prompts/restore', methods=['POST'])
def restore_prompt_version():
    """Restore a historical version of a prompt as the current custom version"""
    try:
        data = request.get_json()
        prompt_name = data.get('prompt_name')
        version = data.get('version')

        if not prompt_name or version is None:
            return jsonify({'error': 'prompt_name and version are required'}), 400

        history = load_prompt_history()
        versions = history.get(prompt_name, [])
        match = next((v for v in versions if v['version'] == version), None)
        if not match:
            return jsonify({'error': f'Version {version} not found for prompt "{prompt_name}"'}), 404

        # Load existing custom prompts
        custom_prompts_path = os.path.join(os.path.dirname(__file__), 'custom_prompts.json')
        custom_prompts = {}
        if os.path.exists(custom_prompts_path):
            with open(custom_prompts_path, 'r', encoding='utf-8') as f:
                custom_prompts = json.load(f)

        # Archive current before restoring
        append_to_history(prompt_name, match['template'], custom_prompts)

        custom_prompts[prompt_name] = {
            'template': match['template'],
            'saved_at': datetime.now().isoformat(),
            'restored_from_version': version
        }
        with open(custom_prompts_path, 'w', encoding='utf-8') as f:
            json.dump(custom_prompts, f, indent=2, ensure_ascii=False)

        return jsonify({'success': True, 'message': f'Restored version {version} of "{prompt_name}"'}), 200
    except Exception as e:
        print(f"Error restoring prompt version: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/prompts/reset', methods=['POST'])
def reset_custom_prompt():
    """Reset a custom prompt to default"""
    try:
        data = request.get_json()
        prompt_name = data.get('prompt_name')
        
        if not prompt_name:
            return jsonify({'error': 'prompt_name is required'}), 400
        
        # Load existing custom prompts
        custom_prompts_path = os.path.join(os.path.dirname(__file__), 'custom_prompts.json')
        custom_prompts = {}
        if os.path.exists(custom_prompts_path):
            with open(custom_prompts_path, 'r', encoding='utf-8') as f:
                custom_prompts = json.load(f)
        
        # Remove prompt
        if prompt_name in custom_prompts:
            del custom_prompts[prompt_name]
            
            with open(custom_prompts_path, 'w', encoding='utf-8') as f:
                json.dump(custom_prompts, f, indent=2, ensure_ascii=False)
        
        return jsonify({'success': True, 'message': f'Prompt "{prompt_name}" reset to default'}), 200
    except Exception as e:
        print(f"Error resetting custom prompt: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/courses', methods=['POST'])
def get_courses():
    """Fetch all active courses from Canvas API (handles pagination)."""
    try:
        data = request.get_json()
        api_token = data.get('api_token')
        canvas_url = data.get('canvas_url', CANVAS_BASE_URL)
        
        if not api_token:
            return jsonify({'error': 'API token is required'}), 400
        
        # Canvas API endpoint for courses
        url = f"{canvas_url}/api/v1/courses"
        headers = {
            'Authorization': f'Bearer {api_token}'
        }
        
        params = {
            'enrollment_state': 'active',
            'per_page': 100,  # Fetch up to 100 courses
            'include[]': 'term'
        }
        
        response = requests.get(url, headers=headers, params=params, timeout=20)
        
        if response.status_code == 200:
            courses = []
            seen_ids = set()
            while True:
                batch = response.json() or []
                for course in batch:
                    cid = course.get('id')
                    if cid in seen_ids:
                        continue
                    seen_ids.add(cid)
                    courses.append(course)
                next_link = response.links.get('next', {}).get('url')
                if not next_link:
                    break
                response = requests.get(next_link, headers=headers, timeout=20)
                if response.status_code != 200:
                    return jsonify({
                        'error': f'Canvas API pagination error: {response.status_code}',
                        'message': response.text
                    }), response.status_code

            # Extract course id and name
            course_list = [
                {
                    'id': course.get('id'),
                    'name': course.get('name'),
                    'course_code': course.get('course_code'),
                    'term_name': (course.get('term') or {}).get('name', '')
                }
                for course in courses
                if course.get('name')  # Only include courses with names
            ]
            return jsonify({'courses': course_list}), 200
        else:
            return jsonify({
                'error': f'Canvas API error: {response.status_code}',
                'message': response.text
            }), response.status_code
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/canvas-user', methods=['POST'])
def get_canvas_user():
    """Fetch current Canvas user profile from API token"""
    try:
        data = request.get_json()
        api_token = data.get('api_token')
        canvas_url = data.get('canvas_url', CANVAS_BASE_URL)

        if not api_token:
            return jsonify({'error': 'API token is required'}), 400

        url = f"{canvas_url}/api/v1/users/self/profile"
        headers = {
            'Authorization': f'Bearer {api_token}'
        }

        response = requests.get(url, headers=headers)
        if response.status_code != 200:
            return jsonify({
                'error': f'Canvas API error: {response.status_code}',
                'message': response.text
            }), response.status_code

        profile = response.json()
        return jsonify({
            'name': profile.get('name', ''),
            'short_name': profile.get('short_name', ''),
            'login_id': profile.get('login_id', '')
        }), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/course-instructors', methods=['POST'])
def get_course_instructors():
    """Fetch actual instructor(s) assigned to a specific Canvas course."""
    try:
        data = request.get_json()
        course_id = data.get('course_id')
        api_token = data.get('api_token')
        canvas_url = data.get('canvas_url', CANVAS_BASE_URL)

        if not course_id or not api_token:
            return jsonify({'error': 'course_id and api_token are required'}), 400

        headers = {'Authorization': f'Bearer {api_token}'}
        url = f"{canvas_url}/api/v1/courses/{course_id}/users"
        params = [('enrollment_type[]', 'teacher'), ('per_page', 100)]
        response = requests.get(url, headers=headers, params=params, timeout=20)

        if response.status_code != 200:
            return jsonify({
                'error': f'Canvas API error: {response.status_code}',
                'message': response.text
            }), response.status_code

        users = response.json() or []
        seen = set()
        instructors = []
        for u in users:
            uid = u.get('id')
            if uid in seen:
                continue
            seen.add(uid)
            name = u.get('name') or u.get('sortable_name') or u.get('short_name') or ''
            if name:
                instructors.append({'id': uid, 'name': name})

        return jsonify({
            'success': True,
            'course_id': course_id,
            'instructors': instructors
        }), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/extract-course', methods=['POST'])
def extract_course():
    """Extract materials from a single course (all or specific type)"""
    try:
        data = request.get_json()
        course_id = data.get('course_id')
        api_token = data.get('api_token')
        canvas_url = data.get('canvas_url', CANVAS_BASE_URL)
        force_refresh = data.get('force_refresh', False)  # Option to force re-extraction
        material_type = data.get('material_type', 'all')  # NEW: specific material type or 'all'
        
        if not course_id or not api_token:
            return jsonify({'error': 'Course ID and API token are required'}), 400
        
        # Check if course data already exists
        course_dir = os.path.join(DATA_DIR, f"course_{course_id}")
        metadata_file = os.path.join(course_dir, "metadata.json")
        
        if os.path.exists(metadata_file) and not force_refresh:
            print(f"\n=== Course {course_id} already extracted ===")
            print(f"Loading existing data from {course_dir}")
            
            # Load and return existing data
            try:
                existing_data = load_course_data(course_id)
                if existing_data:
                    cached_discussions = existing_data.get('discussions', [])
                    if not cached_discussions and existing_data.get('discussions_combined'):
                        cached_discussions = [
                            {
                                'id': d.get('id'),
                                'name': d.get('name', 'Discussion'),
                                'type': 'discussion',
                                'content': '',
                                'source_type': 'discussion_prompt'
                            }
                            for d in existing_data['discussions_combined'].get('discussion_index', [])
                        ]
                    return jsonify({
                        'success': True,
                        'data': {
                            'metadata': existing_data['metadata'],
                            'materials': {
                                'syllabus': existing_data.get('syllabus', []),
                                'lectures': existing_data.get('files', []),
                                'assignments': existing_data.get('assignments', []),
                                'quizzes': existing_data.get('quizzes', []),
                                'quiz_questions': existing_data.get('quiz_questions', []),
                                'labs': [],
                                'exams': [],
                                'discussions': cached_discussions
                            },
                            'statistics': {
                                'total_materials': (
                                    len(existing_data.get('syllabus', [])) +
                                    len(existing_data.get('files', [])) +
                                    len(existing_data.get('assignments', [])) +
                                    len(existing_data.get('quizzes', [])) +
                                    len(existing_data.get('quiz_questions', [])) +
                                    len(cached_discussions)
                                ),
                                'extraction_errors': []
                            }
                        },
                        'message': f'Loaded existing data (extracted: {existing_data["metadata"].get("extracted_at", "unknown")})',
                        'from_cache': True
                    }), 200
            except Exception as e:
                print(f"Error loading existing data: {e}")
                print("Proceeding with fresh extraction...")

        ai_key = resolve_ai_api_key(data)
        llm_model = resolve_llm_model(data)
        needs_gemini = material_type in ('all', 'syllabus', 'lectures', 'all_selected')
        if needs_gemini and not ai_key:
            return jsonify({
                'error': 'AI API key is required for syllabus, lecture, or combined material extraction. Add your key under Configure Canvas, or set AI_API_KEY / GEMINI_API_KEY on the server.'
            }), 400
        if needs_gemini and not llm_model:
            return jsonify({
                'error': 'LLM model name is required for syllabus, lecture, or combined material extraction. Enter it under Configure Canvas, or set LLM_MODEL on the server.'
            }), 400
        
        print(f"\n=== Extracting course {course_id} ===")
        
        headers = {
            'Authorization': f'Bearer {api_token}'
        }
        
        # Initialize extracted data structure
        extracted_data = {
            'metadata': {
                'course_id': course_id,
                'course_name': '',
                'course_code': '',
                'course_term': '',
                'extracted_at': datetime.now().isoformat(),
                'canvas_url': canvas_url
            },
            'materials': {
                'syllabus': [],
                'lectures': [],
                'assignments': [],
                'quizzes': [],
                'quiz_questions': [],
                'labs': [],
                'exams': [],
                'discussions': []
            },
            'statistics': {
                'total_materials': 0,
                'extraction_errors': []
            }
        }
        
        print(f"\n=== Starting extraction for course {course_id} ===")
        
        # Get course info
        try:
            print(f"Fetching course info...")
            course_info = fetch_course_info(course_id, canvas_url, headers)
            if course_info:
                extracted_data['metadata']['course_name'] = course_info.get('name', '')
                extracted_data['metadata']['course_code'] = course_info.get('course_code', '')
                term = course_info.get('term') or {}
                extracted_data['metadata']['course_term'] = term.get('name', '')
                print(f"Course: {extracted_data['metadata']['course_name']} ({extracted_data['metadata']['course_term']})")
        except Exception as e:
            error_msg = f"Error fetching course info: {str(e)}"
            print(error_msg)
            extracted_data['statistics']['extraction_errors'].append(error_msg)
        
        # Extract syllabus (if requested)
        if material_type in ['all', 'syllabus']:
            try:
                print(f"Extracting syllabus...")
                syllabus = extract_syllabus(course_id, canvas_url, headers, api_key=ai_key, llm_model=llm_model)
                if syllabus:
                    extracted_data['materials']['syllabus'].append(syllabus)
                    print(f"✓ Syllabus extracted and filtered")
            except Exception as e:
                error_msg = f"Error extracting syllabus: {str(e)}"
                print(error_msg)
                extracted_data['statistics']['extraction_errors'].append(error_msg)
        
        # Extract assignments (ONLY descriptions/instructions, NOT submissions)
        if material_type in ['all', 'assignments', 'graded_assignments', 'all_selected']:
            try:
                selected_assignment_ids = data.get('selected_assignment_ids') if material_type in ['graded_assignments', 'all_selected'] else None
                if material_type in ['graded_assignments', 'all_selected']:
                    print(f"Extracting selected graded assignments...")
                else:
                    print(f"Extracting assignments...")
                assignments = extract_assignments(course_id, canvas_url, headers, assignment_ids=selected_assignment_ids)
                categorize_materials(assignments, extracted_data['materials'])
                print(f"✓ Extracted {len(assignments)} assignments")
            except Exception as e:
                error_msg = f"Error extracting assignments: {str(e)}"
                print(error_msg)
                extracted_data['statistics']['extraction_errors'].append(error_msg)
        
        extracted_data['material_type'] = material_type
        
        # Skip pages extraction (not needed for comparison)
        print(f"Skipping pages extraction (not needed)")
        
        # Extract quizzes (both quiz overviews AND individual questions)
        if material_type in ['all', 'quizzes']:
            try:
                print(f"Extracting quizzes...")
                quizzes, quiz_questions = extract_quizzes(course_id, canvas_url, headers)
                categorize_materials(quizzes, extracted_data['materials'])
                categorize_materials(quiz_questions, extracted_data['materials'])  # Store questions separately
                print(f"✓ Extracted {len(quizzes)} quizzes with {len(quiz_questions)} individual questions")
            except Exception as e:
                error_msg = f"Error extracting quizzes: {str(e)}"
                print(error_msg)
                extracted_data['statistics']['extraction_errors'].append(error_msg)
        
        # Extract discussions
        if material_type in ['all', 'discussions', 'all_selected']:
            try:
                selected_discussion_ids = data.get('selected_discussion_ids') if material_type == 'discussions' else None
                print(f"Extracting discussions...")
                discussions = extract_discussions(course_id, canvas_url, headers, selected_discussion_ids=selected_discussion_ids)
                categorize_materials(discussions, extracted_data['materials'])
                print(f"✓ Extracted {len(discussions)} discussions")
            except Exception as e:
                error_msg = f"Error extracting discussions: {str(e)}"
                print(error_msg)
                extracted_data['statistics']['extraction_errors'].append(error_msg)
        
        # Extract files (PDFs, documents) - Text extraction only
        if material_type in ['all', 'lectures', 'all_selected']:
            try:
                lectures_only = (material_type in ['lectures', 'all_selected'])
                selected_file_ids = data.get('selected_file_ids', [])  # NEW: Selected file IDs
                
                if lectures_only:
                    if selected_file_ids:
                        print(f"Extracting {len(selected_file_ids)} selected lecture files...")
                    else:
                        print(f"Extracting lecture materials (ungraded files: PDFs, PPTs, Word docs)...")
                else:
                    print(f"Extracting text from files...")
                    
                api_key = ai_key if lectures_only else None
                lm = llm_model if lectures_only else None
                files = extract_files(
                    course_id, canvas_url, headers,
                    lectures_only=lectures_only, selected_file_ids=selected_file_ids,
                    api_key=api_key, llm_model=lm
                )
                categorize_materials(files, extracted_data['materials'])
                print(f"✓ Extracted {len(files)} files")
                if lectures_only:
                    lec_count = len(extracted_data['materials'].get('lectures', []))
                    print(f"  → {lec_count} files categorized as lectures")
            except Exception as e:
                error_msg = f"Error extracting files: {str(e)}"
                print(error_msg)
                extracted_data['statistics']['extraction_errors'].append(error_msg)
        
        # Calculate statistics
        extracted_data['statistics']['total_materials'] = sum([
            len(extracted_data['materials']['syllabus']),
            len(extracted_data['materials']['lectures']),
            len(extracted_data['materials']['assignments']),
            len(extracted_data['materials']['quizzes']),
            len(extracted_data['materials']['quiz_questions']),
            len(extracted_data['materials']['labs']),
            len(extracted_data['materials']['exams']),
            len(extracted_data['materials']['discussions'])
        ])
        
        print(f"=== Extraction complete: {extracted_data['statistics']['total_materials']} materials ===\n")
        
        # Save to JSON file
        try:
            save_course_data(course_id, extracted_data)
        except Exception as e:
            print(f"Warning: Could not save to file: {e}")
        
        return jsonify({
            'success': True,
            'data': extracted_data,
            'message': f'Successfully extracted {extracted_data["statistics"]["total_materials"]} materials'
        }), 200
        
    except Exception as e:
        error_msg = str(e)
        print(f"FATAL ERROR in extract_course: {error_msg}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': error_msg}), 500

def fetch_course_info(course_id, canvas_url, headers):
    """Fetch basic course information including enrollment term"""
    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}"
        response = requests.get(url, headers=headers, params={'include[]': 'term'})
        if response.status_code == 200:
            return response.json()
    except Exception as e:
        print(f"Error fetching course info: {e}")
    return None


def parse_syllabus_into_sections(content, format='html'):
    """
    Parse syllabus content into a dict of section_heading -> content.
    format: 'html' (Canvas syllabus_body) or 'text' (e.g. from PDF).
    Section headers are detected from HTML h1/h2/h3/h4 or from bold/strong;
    for plain text, from short lines, lines ending with colon, or ALL CAPS lines.
    """
    sections = {}
    if not content or not content.strip():
        return sections

    if format == 'html':
        try:
            soup = BeautifulSoup(content, 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            current_header = None
            current_parts = []
            # Walk all meaningful elements
            for el in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'div', 'li', 'span']):
                text = el.get_text(separator=' ', strip=True)
                if not text:
                    continue
                is_header = False
                if el.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    is_header = True
                elif el.name == 'p' and len(text) < 120:
                    # Short paragraph might be a section title (e.g. <p><strong>Course Description</strong></p>
                    if el.find(['strong', 'b']) and len(list(el.stripped_strings)) <= 2:
                        is_header = True
                if is_header:
                    if current_header is not None and current_parts:
                        sections[current_header] = '\n'.join(current_parts).strip()
                    current_header = text
                    current_parts = []
                else:
                    current_parts.append(text)
            if current_header is not None and current_parts:
                sections[current_header] = '\n'.join(current_parts).strip()
            if not sections and content.strip():
                sections['Content'] = clean_html(content)
        except Exception as e:
            print(f"  Warning: could not parse syllabus HTML into sections: {e}")
            sections['Content'] = clean_html(content)
        return sections

    # Plain text (e.g. from PDF)
    lines = content.splitlines()
    current_header = 'Introduction'
    current_parts = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        # Heuristic: header = short line, or ends with colon, or mostly caps
        is_likely_header = (
            len(stripped) < 80 and
            (stripped.endswith(':') or stripped.isupper() or
             (len(stripped.split()) <= 8 and stripped[0].isupper()))
        )
        if is_likely_header and len(current_parts) > 0:
            sections[current_header] = '\n'.join(current_parts).strip()
            current_header = stripped.rstrip(':')
            current_parts = []
        else:
            current_parts.append(stripped)
    if current_parts:
        sections[current_header] = '\n'.join(current_parts).strip()
    if not sections and content.strip():
        sections['Content'] = content.strip()
    return sections


def filter_syllabus_with_gemini(raw_syllabus_text, api_key, llm_model):
    """Use LLM to extract only Course Description, Learning Outcomes, Topics/Schedule. Ignores all policy and boilerplate. Input may be raw unsectioned text."""
    try:
        from google import genai
        from prompt_loader import load_prompt

        client = genai.Client(api_key=api_key)
        raw_str = (raw_syllabus_text or "").strip()
        if not raw_str:
            return {"course_description": "", "learning_outcomes": "", "topics_schedule": ""}

        prompt_data = load_prompt('syllabus_filtering')
        template = prompt_data['template']
        prompt = template.replace('{raw_syllabus_text}', raw_str)
        prompt = prompt.replace('{{', '{').replace('}}', '}')

        response = client.models.generate_content(
            model=llm_model,
            contents=prompt
        )
        
        # Extract JSON from response
        response_text = response.text.strip()
        
        # Remove markdown code blocks if present
        if '```json' in response_text:
            response_text = response_text.split('```json')[1].split('```')[0].strip()
        elif '```' in response_text:
            response_text = response_text.split('```')[1].split('```')[0].strip()
        
        # Parse JSON
        filtered_data = json.loads(response_text)
        
        print(f"✓ LLM filtered syllabus successfully")
        print(f"  - Course Description: {len(filtered_data.get('course_description', '')) > 0}")
        print(f"  - Learning Outcomes: {len(filtered_data.get('learning_outcomes', '')) > 0}")
        print(f"  - Topics/Schedule: {len(filtered_data.get('topics_schedule', '')) > 0}")
        
        return filtered_data
        
    except Exception as e:
        print(f"Error filtering syllabus with LLM: {e}")
        import traceback
        traceback.print_exc()
        # Return empty structure if filtering fails
        return {
            "course_description": "",
            "learning_outcomes": "",
            "topics_schedule": "",
            "error": str(e)
        }

def extract_syllabus(course_id, canvas_url, headers, api_key=None, llm_model=None):
    """Extract course syllabus from Canvas or PDF, then LLM extracts only description, outcomes, and schedule (ignores all policy)."""
    raw_content = None
    syllabus_source = None

    try:
        # STEP 1: Try to get syllabus from syllabus_body field
        url = f"{canvas_url}/api/v1/courses/{course_id}"
        params = {'include': ['syllabus_body']}
        response = requests.get(url, headers=headers, params=params)

        if response.status_code == 200:
            course_data = response.json()
            syllabus_body = course_data.get('syllabus_body', '')

            if syllabus_body and syllabus_body.strip():
                cleaned_text = clean_html(syllabus_body)

                # Check if it's meaningful content (not just a PDF reference)
                has_pdf_reference = '.pdf' in cleaned_text.lower() or 'pdf' in cleaned_text.lower()
                is_just_reference = len(cleaned_text) < 500 or 'linked above' in cleaned_text.lower() or 'see attached' in cleaned_text.lower()

                if len(cleaned_text) > 500 and not has_pdf_reference:
                    raw_content = cleaned_text
                    syllabus_source = 'syllabus_body'
                    print(f"✓ Extracted syllabus from Canvas syllabus_body field")
                elif len(cleaned_text) > 100 and not is_just_reference and not has_pdf_reference:
                    raw_content = cleaned_text
                    syllabus_source = 'syllabus_body'
                    print(f"✓ Extracted syllabus from Canvas syllabus_body field")
        
        # STEP 2: If no good content from syllabus_body, search for syllabus PDF in files
        if not raw_content:
            print(f"  Syllabus body is empty or just a PDF reference, searching course files...")
            
            # Try to extract PDF filename from syllabus_body if it's mentioned
            pdf_filename_in_body = None
            if syllabus_body:
                # Look for .pdf filename in the body text
                import re
                pdf_match = re.search(r'([A-Za-z0-9_\-\.]+\.pdf)', syllabus_body)
                if pdf_match:
                    pdf_filename_in_body = pdf_match.group(1)
                    print(f"  Found PDF reference in body: {pdf_filename_in_body}")
            
            # Fetch all course files
            files_url = f"{canvas_url}/api/v1/courses/{course_id}/files"
            params = {'per_page': 100}
            files_response = requests.get(files_url, headers=headers, params=params)
            
            if files_response.status_code == 200:
                files = files_response.json()
                
                print(f"  Found {len(files)} files in course")
                
                # Debug: Show all PDF files
                pdf_files = [f.get('filename', '') for f in files if f.get('filename', '').lower().endswith('.pdf')]
                if pdf_files:
                    print(f"  PDF files found: {pdf_files[:5]}")  # Show first 5
                
                # Look for syllabus PDF (case-insensitive)
                syllabus_file = None
                
                # First, try to find the exact file mentioned in syllabus_body
                if pdf_filename_in_body:
                    for file in files:
                        if file.get('filename', '').lower() == pdf_filename_in_body.lower():
                            syllabus_file = file
                            print(f"  Found exact match: {file.get('filename')}")
                            break
                
                # If not found, look for any file with "syllabus" in name
                if not syllabus_file:
                    for file in files:
                        filename = file.get('filename', '').lower()
                        display_name = file.get('display_name', '').lower()
                        
                        # Check both filename and display_name for "syllabus"
                        if 'syllabus' in filename and filename.endswith('.pdf'):
                            syllabus_file = file
                            break
                        elif 'syllabus' in display_name and filename.endswith('.pdf'):
                            syllabus_file = file
                            break
                
                if syllabus_file:
                    print(f"  Found syllabus PDF: {syllabus_file['filename']}")
                    
                    # Download and extract PDF text
                    file_url = syllabus_file.get('url')
                    if file_url:
                        try:
                            file_response = requests.get(file_url, headers=headers)
                            if file_response.status_code == 200:
                                # Save temporarily and extract
                                import tempfile
                                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                                    tmp_file.write(file_response.content)
                                    tmp_path = tmp_file.name
                                
                                # Extract PDF text
                                pdf_text = extract_pdf_text(tmp_path)
                                
                                # Clean up temp file
                                import os
                                os.unlink(tmp_path)
                                
                                if pdf_text:
                                    raw_content = pdf_text
                                    syllabus_source = f"pdf_file:{syllabus_file['filename']}"
                                    print(f"  ✓ Extracted {len(pdf_text)} characters from syllabus PDF")
                                else:
                                    print(f"  ⚠ Could not extract text from PDF")
                        except Exception as e:
                            print(f"  Error downloading/extracting PDF: {e}")
                else:
                    print(f"  ⚠ No file with 'syllabus' in name found")
                    # Try broader search - look for any PDF with course code or common syllabus patterns
                    for file in files:
                        filename = file.get('filename', '')
                        display_name = file.get('display_name', '')
                        
                        # Look for patterns like: COT4210, course syllabus, etc.
                        if filename.lower().endswith('.pdf'):
                            # Check if it might be a syllabus based on common patterns
                            name_lower = (filename + ' ' + display_name).lower()
                            if any(pattern in name_lower for pattern in ['course', 'outline', 'overview', str(course_id)]):
                                print(f"  Trying potential syllabus file: {filename}")
                                syllabus_file = file
                                break
                    
                    if not syllabus_file:
                        print(f"  ⚠ No syllabus PDF found in course files")
        
        # STEP 3: If we have content, send raw text to Gemini (no sectioning required)
        if raw_content and len(raw_content.strip()) > 50:
            print(f"✓ Extracted RAW syllabus ({len(raw_content)} chars) from {syllabus_source}")

            if api_key and llm_model:
                print(f"🤖 Sending raw syllabus to LLM (extract only description, outcomes, schedule; ignore all policy)...")
                filtered_data = filter_syllabus_with_gemini(raw_content, api_key, llm_model)

                trimmed_syllabus = {
                    'id': f'syllabus_{course_id}',
                    'name': 'Course Syllabus (Filtered)',
                    'type': 'syllabus',
                    'filtered_data': filtered_data,
                    'content': json.dumps(filtered_data, indent=2),
                    'graded': False,
                    'points': 0,
                    'source_type': 'syllabus_filtered',
                    'source': syllabus_source,
                    'raw_content': raw_content,
                }
                print(f"✓ Syllabus filtered and ready for comparison")
                return trimmed_syllabus
            else:
                raw_syllabus = {
                    'id': f'syllabus_raw_{course_id}',
                    'name': 'Course Syllabus (Raw)',
                    'type': 'syllabus_raw',
                    'content': raw_content,
                    'graded': False,
                    'points': 0,
                    'source_type': 'syllabus_raw_untrimmed',
                    'source': syllabus_source
                }
                print(f"⚠ No API key provided - returning raw syllabus (not filtered)")
                return raw_syllabus
        else:
            print(f"  ⚠ No meaningful syllabus content found")
            
    except Exception as e:
        print(f"Error extracting syllabus: {e}")
        import traceback
        traceback.print_exc()
    
    return None

def is_graded(assignment):
    """Determine if an assignment is graded using Canvas API fields."""
    if assignment.get('grading_type') == 'not_graded':
        return False
    if assignment.get('points_possible', 0) == 0:
        return False
    return True


def _file_ids_from_description(description_html):
    """Extract Canvas file IDs from assignment description HTML (e.g. /files/123/download)."""
    if not description_html:
        return []
    # Match /files/123/ or /files/123/download or /courses/.../files/123/...
    ids = set()
    for m in re.finditer(r'/files/(\d+)(?:/|$)', description_html):
        ids.add(int(m.group(1)))
    return list(ids)


def _graded_file_ids_for_course(course_id, canvas_url, headers):
    """Return set of file IDs that are linked from graded assignments (same logic as tests).
    Used to exclude those files from the lecture materials list (show only non-graded files)."""
    graded_file_ids = set()
    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}/assignments"
        r = requests.get(url, headers=headers, params={'per_page': 100}, timeout=10)
        if r.status_code != 200:
            return graded_file_ids
        for assignment in r.json():
            if not is_graded(assignment):
                continue
            desc = assignment.get('description') or ''
            for fid in _file_ids_from_description(desc):
                graded_file_ids.add(fid)
    except Exception as e:
        print(f"  Could not fetch graded assignment file IDs: {e}")
    return graded_file_ids


def _fetch_quiz_questions_text(course_id, quiz_id, assignment_name, canvas_url, headers):
    """Fetch quiz questions for an assignment that is a quiz. Returns formatted text."""
    parts = []
    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}/quizzes/{quiz_id}/questions"
        r = requests.get(url, headers=headers, params={'per_page': 100}, timeout=15)
        if r.status_code != 200:
            return ""
        questions = r.json()
        if not questions:
            return ""
        parts.append("\n\n=== QUIZ QUESTIONS (assignment is a quiz) ===\n")
        for idx, q in enumerate(questions, 1):
            q_text = clean_html(q.get('question_text', ''))
            q_type = q.get('question_type', 'unknown')
            pts = q.get('points_possible', 0)
            parts.append(f"\nQuestion {idx} ({q_type}, {pts} pts):\n{q_text}")
            answers = q.get('answers', [])
            if answers and q_type in ['multiple_choice_question', 'true_false_question', 'multiple_answers_question']:
                for a in answers:
                    at = clean_html(a.get('text', ''))
                    if at:
                        parts.append(f"  - {at}")
        return "\n".join(parts)
    except Exception as e:
        print(f"  Could not fetch quiz questions for quiz {quiz_id}: {e}")
        return ""


def _fetch_and_extract_file_content(course_id, file_id, canvas_url, headers):
    """Download a course file by ID and extract text (PDF, docx, pptx). Returns string or None."""
    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}/files/{file_id}"
        r = requests.get(url, headers=headers, timeout=15)
        if r.status_code != 200:
            return None
        info = r.json()
        file_url = info.get('url')
        display_name = (info.get('display_name') or '').lower()
        if not file_url:
            return None
        # Download file content
        file_r = requests.get(file_url, headers=headers, timeout=30)
        if file_r.status_code != 200:
            return None
        content = None
        if display_name.endswith('.pdf'):
            content = extract_pdf_content(file_url, headers)
        elif display_name.endswith('.pptx'):
            content = extract_pptx_content(file_url, headers)
        elif display_name.endswith('.docx'):
            content = extract_docx_content(file_url, headers)
        return content
    except Exception as e:
        print(f"  Could not extract file {file_id}: {e}")
        return None


def extract_assignments(course_id, canvas_url, headers, assignment_ids=None):
    """Extract selected assignments for comparison. For each assignment we check two things:
    1) Is it uploaded as a quiz? -> extract quiz questions.
    2) Does it have a PDF or any linked file (e.g. in description)? -> extract file content.
    We always include the written description when present.
    Result: quiz questions and/or file content and/or description are combined per assignment.
    If assignment_ids is provided, only those assignments are included. Multiple selected -> all are combined into one JSON per course (via combine_graded_assignments)."""
    materials = []
    
    skip_keywords = ['survey', 'participation', 'attendance', 'pitch', 'introductions', 'introduction', 'flyer']
    
    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}/assignments"
        params = {'per_page': 100}
        response = requests.get(url, headers=headers, params=params)
        
        if response.status_code == 200:
            assignments = response.json()
            for assignment in assignments:
                if assignment_ids is not None and assignment.get('id') not in assignment_ids:
                    continue
                name = assignment.get('name', 'Untitled Assignment')
                
                name_lower = name.lower()
                if any(keyword in name_lower for keyword in skip_keywords):
                    print(f"  ⊗ Skipping: {name} (survey/participation/attendance/pitch/flyer)")
                    continue
                
                description = assignment.get('description', '')
                content_parts = []
                quiz_text = None
                
                # 1) Written description (always)
                desc_text = clean_html(description)
                if desc_text.strip():
                    content_parts.append(desc_text)
                
                # 2) If assignment is a quiz (Canvas quiz_id present), extract questions
                quiz_id = assignment.get('quiz_id')
                if quiz_id:
                    print(f"  Assignment is a quiz: {name}, fetching questions...")
                    quiz_text = _fetch_quiz_questions_text(course_id, quiz_id, name, canvas_url, headers)
                    if quiz_text:
                        content_parts.append(quiz_text)
                
                # 3) If description links to PDF or other files, extract their content
                file_ids = _file_ids_from_description(description)
                for fid in file_ids:
                    print(f"  Extracting linked file id={fid} for assignment: {name}")
                    file_content = _fetch_and_extract_file_content(course_id, fid, canvas_url, headers)
                    if file_content and file_content.strip():
                        content_parts.append(f"\n\n=== CONTENT FROM ATTACHED/LINKED FILE (id={fid}) ===\n\n{file_content}")
                
                content = "\n\n".join(content_parts) if content_parts else "(no content extracted)"
                
                # Explicit content_sources: what we checked and extracted (quiz, files, or just description)
                content_sources = []
                if desc_text.strip():
                    content_sources.append('description')
                if quiz_id and quiz_text:
                    content_sources.append('quiz_questions')
                if file_ids:
                    content_sources.append('linked_files')
                if not content_sources:
                    content_sources.append('none')
                
                material = {
                    'id': assignment.get('id'),
                    'name': name,
                    'type': 'assignment',
                    'content': content,
                    'content_sources': content_sources,
                    'is_quiz': bool(quiz_id),
                    'had_linked_files': len(file_ids) > 0,
                    'graded': is_graded(assignment),
                    'points': assignment.get('points_possible', 0),
                    'due_date': assignment.get('due_at'),
                    'position': assignment.get('position', 999),
                    'source_type': 'assignment_description',
                    'submission_types': assignment.get('submission_types', []),
                    'quiz_id': quiz_id,
                    'linked_file_ids': file_ids
                }
                materials.append(material)
            
            materials.sort(key=lambda x: x['position'])
            
    except Exception as e:
        print(f"Error extracting assignments: {e}")
    
    return materials

def extract_pages(course_id, canvas_url, headers):
    """Extract course pages (lecture notes, content pages)"""
    materials = []
    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}/pages"
        params = {'per_page': 100, 'sort': 'created_at'}
        response = requests.get(url, headers=headers, params=params)
        
        if response.status_code == 200:
            pages = response.json()
            for idx, page in enumerate(pages):
                page_url = page.get('url')
                # Fetch full page content
                page_detail_url = f"{canvas_url}/api/v1/courses/{course_id}/pages/{page_url}"
                page_response = requests.get(page_detail_url, headers=headers)
                
                if page_response.status_code == 200:
                    page_data = page_response.json()
                    body = page_data.get('body', '')
                    content = clean_html(body)
                    
                    material = {
                        'id': page_data.get('page_id'),
                        'name': page_data.get('title', 'Untitled Page'),
                        'type': 'page',
                        'content': content,
                        'graded': False,
                        'points': 0,
                        'position': idx,
                        'created_at': page_data.get('created_at'),
                        'source_type': 'page_content'
                    }
                    
                    materials.append(material)
    except Exception as e:
        print(f"Error extracting pages: {e}")
    
    return materials

def extract_quizzes(course_id, canvas_url, headers):
    """Extract quiz descriptions AND questions (NOT student answers)"""
    materials = []
    quiz_questions = []  # Separate list for individual quiz questions
    
    # Keywords to skip (surveys, participation, attendance, flyers, etc.)
    skip_keywords = ['survey', 'participation', 'attendance', 'pitch', 'introductions', 'introduction', 'flyer']
    
    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}/quizzes"
        params = {'per_page': 100}
        response = requests.get(url, headers=headers, params=params)
        
        if response.status_code == 200:
            quizzes = response.json()
            for quiz in quizzes:
                quiz_id = quiz.get('id')
                quiz_name = quiz.get('title', 'Untitled Quiz')
                
                # Skip quizzes with unwanted keywords (surveys, attendance, flyers, etc.)
                quiz_name_lower = quiz_name.lower()
                if any(keyword in quiz_name_lower for keyword in skip_keywords):
                    print(f"  ⊗ Skipping quiz: {quiz_name} (survey/participation/attendance/pitch/flyer)")
                    continue
                
                description = quiz.get('description', '')
                content_parts = [clean_html(description)] if description else []
                
                # Fetch quiz questions
                try:
                    questions_url = f"{canvas_url}/api/v1/courses/{course_id}/quizzes/{quiz_id}/questions"
                    questions_response = requests.get(questions_url, headers=headers, params={'per_page': 100}, timeout=10)
                    
                    if questions_response.status_code == 200:
                        questions = questions_response.json()
                        
                        if questions:
                            content_parts.append("\n\n=== QUIZ QUESTIONS ===\n")
                            
                            for idx, question in enumerate(questions, 1):
                                q_text = clean_html(question.get('question_text', ''))
                                q_type = question.get('question_type', 'unknown')
                                q_points = question.get('points_possible', 0)
                                
                                # Add to quiz content
                                content_parts.append(f"\nQuestion {idx} ({q_type}, {q_points} pts):")
                                content_parts.append(q_text)
                                
                                # Add answer choices for multiple choice/true-false
                                answers = question.get('answers', [])
                                choices_text = []
                                if answers and q_type in ['multiple_choice_question', 'true_false_question', 'multiple_answers_question']:
                                    content_parts.append("Choices:")
                                    for answer in answers:
                                        answer_text = clean_html(answer.get('text', ''))
                                        if answer_text:
                                            content_parts.append(f"  - {answer_text}")
                                            choices_text.append(answer_text)
                                
                                # CREATE INDIVIDUAL QUIZ QUESTION ENTRY
                                question_content = f"Question Type: {q_type}\nPoints: {q_points}\n\nQuestion:\n{q_text}"
                                if choices_text:
                                    question_content += "\n\nAnswer Choices:\n" + "\n".join([f"- {c}" for c in choices_text])
                                
                                quiz_questions.append({
                                    'id': f"{quiz_id}_q{question.get('id', idx)}",
                                    'name': f"{quiz_name} - Question {idx}",
                                    'type': 'quiz_question',
                                    'content': question_content,
                                    'graded': quiz.get('points_possible', 0) > 0,
                                    'points': q_points,
                                    'quiz_id': quiz_id,
                                    'quiz_name': quiz_name,
                                    'question_number': idx,
                                    'question_type': q_type,
                                    'source_type': 'quiz_question'
                                })
                            
                            print(f"  ✓ Extracted {len(questions)} questions from quiz: {quiz_name}")
                    
                except Exception as e:
                    print(f"  Could not fetch questions for quiz {quiz_id}: {e}")
                
                # Combine all content for the quiz overview
                content = '\n'.join(content_parts)
                
                material = {
                    'id': quiz_id,
                    'name': quiz.get('title', 'Untitled Quiz'),
                    'type': 'quiz',
                    'content': content,
                    'graded': quiz.get('points_possible', 0) > 0,
                    'points': quiz.get('points_possible', 0),
                    'due_date': quiz.get('due_at'),
                    'position': quiz.get('position', 999),
                    'question_count': quiz.get('question_count', 0),
                    'source_type': 'quiz_with_questions'
                }
                
                materials.append(material)
            
            # Sort by position
            materials.sort(key=lambda x: x['position'])
            quiz_questions.sort(key=lambda x: (x['quiz_id'], x['question_number']))
            
    except Exception as e:
        print(f"Error extracting quizzes: {e}")
    
    # Return both quiz overviews AND individual questions
    return materials, quiz_questions

def _teacher_user_ids(course_id, canvas_url, headers):
    """Get teacher user IDs for a course."""
    teacher_ids = set()
    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}/users"
        params = [('enrollment_type[]', 'teacher'), ('per_page', 100)]
        response = requests.get(url, headers=headers, params=params, timeout=15)
        if response.status_code == 200:
            for user in response.json():
                uid = user.get('id')
                if isinstance(uid, int):
                    teacher_ids.add(uid)
    except Exception as e:
        print(f"Warning: Could not fetch teacher IDs for course {course_id}: {e}")
    return teacher_ids

def list_professor_discussion_topics(course_id, canvas_url, headers):
    """Return discussion topics created by course teachers."""
    url = f"{canvas_url}/api/v1/courses/{course_id}/discussion_topics"
    params = {'per_page': 100}
    response = requests.get(url, headers=headers, params=params, timeout=20)
    if response.status_code != 200:
        return []

    discussions = response.json()
    teacher_ids = _teacher_user_ids(course_id, canvas_url, headers)
    filtered = []
    for discussion in discussions:
        author = discussion.get('author') or {}
        author_id = author.get('id') or discussion.get('author_id') or discussion.get('user_id')
        if teacher_ids and author_id not in teacher_ids:
            continue
        filtered.append(discussion)
    return filtered

def extract_discussions(course_id, canvas_url, headers, selected_discussion_ids=None):
    """Extract professor-created discussion topics (prompt only, not replies)."""
    materials = []
    if selected_discussion_ids is None:
        selected_discussion_ids = []
    selected_discussion_ids = [int(did) for did in selected_discussion_ids if str(did).isdigit()]
    
    # Keywords to skip (introductions, attendance, flyers, etc.)
    skip_keywords = ['survey', 'participation', 'attendance', 'pitch', 'introductions', 'introduction', 'flyer']
    
    try:
        discussions = list_professor_discussion_topics(course_id, canvas_url, headers)
        for discussion in discussions:
            discussion_id = discussion.get('id')
            if selected_discussion_ids and discussion_id not in selected_discussion_ids:
                continue
            title = discussion.get('title', 'Untitled Discussion')
            
            # Skip discussions with unwanted keywords
            title_lower = title.lower()
            if any(keyword in title_lower for keyword in skip_keywords):
                print(f"  ⊗ Skipping discussion: {title} (survey/participation/attendance/pitch/introduction/flyer)")
                continue
            
            message = discussion.get('message', '')
            content = clean_html(message)
            
            material = {
                'id': discussion_id,
                'name': title,
                'type': 'discussion',
                'content': content,
                'graded': discussion.get('assignment') is not None,
                'points': discussion.get('assignment', {}).get('points_possible', 0) if discussion.get('assignment') else 0,
                'source_type': 'discussion_prompt',
                'author_name': (discussion.get('author') or {}).get('display_name', ''),
                'posted_at': discussion.get('posted_at') or discussion.get('created_at')
            }
            
            materials.append(material)
    except Exception as e:
        print(f"Error extracting discussions: {e}")
    
    return materials

def extract_files(course_id, canvas_url, headers, lectures_only=False, selected_file_ids=None, api_key=None, llm_model=None):
    """
    Extract text content from files (PDFs, PowerPoint, Word docs).
    When lectures_only and api_key are set, use OCR flow: file -> base64 images -> LLM OCR -> structured JSON per file; then combine.
    All files are marked as ungraded since Canvas doesn't distinguish.

    Args:
        selected_file_ids: List of file IDs to extract. If None, extract all files.
        api_key: AI/LLM API key for OCR when lectures_only=True.
        llm_model: Model id for OCR when lectures_only=True.
    """
    materials = []

    if selected_file_ids is None:
        selected_file_ids = []
    selected_file_ids = [int(fid) for fid in selected_file_ids]

    try:
        url = f"{canvas_url}/api/v1/courses/{course_id}/files"
        params = {'per_page': 100}
        response = requests.get(url, headers=headers, params=params, timeout=10)

        if response.status_code != 200:
            return materials

        files = response.json()
        use_ocr = lectures_only and api_key and llm_model

        if use_ocr:
            try:
                from google import genai
                from lecture_ocr import process_lecture_file_with_ocr
                client = genai.Client(api_key=api_key)
            except Exception as e:
                print(f"  OCR flow unavailable: {e}")
                use_ocr = False

        for file_info in files:
            if selected_file_ids and file_info.get('id') not in selected_file_ids:
                continue
            try:
                file_name = file_info.get('display_name', '')
                file_url = file_info.get('url')
                mime_type = file_info.get('content-type', '')
                file_size = file_info.get('size', 0)
                file_id = file_info.get('id')

                if file_size > 50 * 1024 * 1024:
                    print(f"Skipping large file: {file_name} ({file_size / (1024*1024):.1f} MB)")
                    continue

                content = None
                source_type = None
                ocr_result = None

                ext = ''
                if file_name.lower().endswith('.pdf') or 'pdf' in (mime_type or '').lower():
                    ext = 'pdf'
                elif file_name.lower().endswith('.pptx'):
                    ext = 'pptx'
                elif file_name.lower().endswith('.docx'):
                    ext = 'docx'
                elif file_name.lower().endswith('.ppt') or file_name.lower().endswith('.doc'):
                    continue
                else:
                    continue

                if use_ocr and ext in ('pdf', 'pptx', 'docx'):
                    print(f"OCR lecture: {file_name}")
                    ocr_result = process_lecture_file_with_ocr(
                        file_url, headers, file_name, file_id, ext, api_key, client=client, llm_model=llm_model
                    )
                    if ocr_result and ocr_result.get('slides'):
                        content = '\n\n'.join(
                            s.get('text', '') or ''
                            for s in ocr_result['slides']
                        )
                        source_type = 'powerpoint' if ext == 'pptx' else 'word' if ext == 'docx' else 'pdf'
                if content is None or (not content.strip() and not ocr_result):
                    # Fallback: text-only extraction
                    if ext == 'pdf':
                        print(f"Extracting PDF (text): {file_name}")
                        content = extract_pdf_content(file_url, headers)
                        source_type = 'pdf'
                    elif ext == 'pptx':
                        print(f"Extracting PowerPoint (text): {file_name}")
                        content = extract_pptx_content(file_url, headers)
                        source_type = 'powerpoint'
                    elif ext == 'docx':
                        print(f"Extracting Word (text): {file_name}")
                        content = extract_docx_content(file_url, headers)
                        source_type = 'word'
                    ocr_result = None

                if content is None:
                    content = ''
                if not content.strip() and not ocr_result:
                    continue

                material = {
                    'id': file_id,
                    'name': file_name,
                    'type': 'file',
                    'content': content.strip() or '(no text extracted)',
                    'graded': False,
                    'points': 0,
                    'source_type': source_type or ext,
                    'file_type': mime_type,
                    'file_size': file_size
                }
                if ocr_result is not None:
                    material['ocr_result'] = ocr_result
                materials.append(material)
                print(f"  ✓ Extracted: {file_name}")
            except Exception as e:
                print(f"Error processing file {file_info.get('display_name', '')}: {e}")
                continue

    except Exception as e:
        print(f"Error extracting files: {e}")

    return materials

def extract_pdf_content(file_url, headers):
    """Extract text from PDF file"""
    try:
        response = requests.get(file_url, headers=headers, timeout=30)
        if response.status_code != 200:
            print(f"Failed to download PDF: {response.status_code}")
            return None
        
        pdf_file = io.BytesIO(response.content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text_content = []
        # Extract ALL pages (no limit)
        for page_num in range(len(pdf_reader.pages)):
            try:
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text:
                    text_content.append(text)
            except Exception as e:
                print(f"  Error on page {page_num}: {e}")
                continue
        
        full_text = '\n'.join(text_content)
        full_text = ' '.join(full_text.split())
        
        return full_text if full_text.strip() else None
        
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        return None

def extract_pdf_text(file_path):
    """Extract text from a local PDF file"""
    try:
        with open(file_path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = []
            # Extract ALL pages
            for page_num in range(len(pdf_reader.pages)):
                try:
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
                except Exception as e:
                    print(f"  Error on page {page_num}: {e}")
                    continue
            
            full_text = '\n'.join(text_content)
            full_text = ' '.join(full_text.split())  # Normalize whitespace
            
            return full_text if full_text.strip() else None
            
    except Exception as e:
        print(f"Error extracting PDF from {file_path}: {e}")
        return None

def extract_pptx_content(file_url, headers):
    """Extract text from PowerPoint (.pptx) file"""
    try:
        response = requests.get(file_url, headers=headers, timeout=30)
        if response.status_code != 200:
            print(f"Failed to download PowerPoint: {response.status_code}")
            return None
        
        pptx_file = io.BytesIO(response.content)
        presentation = Presentation(pptx_file)
        
        text_content = []
        # Extract ALL slides (no limit)
        for slide_num, slide in enumerate(presentation.slides, 1):
            try:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                if slide_text:
                    text_content.append(f"Slide {slide_num}: " + ' '.join(slide_text))
            except Exception as e:
                print(f"  Error on slide {slide_num}: {e}")
                continue
        
        full_text = '\n'.join(text_content)
        full_text = ' '.join(full_text.split())
        
        return full_text if full_text.strip() else None
        
    except Exception as e:
        print(f"Error extracting PowerPoint: {e}")
        return None

def extract_docx_content(file_url, headers):
    """Extract text from Word (.docx) file"""
    try:
        response = requests.get(file_url, headers=headers, timeout=30)
        if response.status_code != 200:
            print(f"Failed to download Word document: {response.status_code}")
            return None
        
        docx_file = io.BytesIO(response.content)
        document = Document(docx_file)
        
        text_content = []
        
        try:
            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
        except Exception as e:
            print(f"  Error extracting paragraphs: {e}")
        
        try:
            for table in document.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content.append(cell.text)
        except Exception as e:
            print(f"  Error extracting tables: {e}")
        
        full_text = '\n'.join(text_content)
        full_text = ' '.join(full_text.split())
        
        return full_text if full_text.strip() else None
        
    except Exception as e:
        print(f"Error extracting Word document: {e}")
        return None

def combine_graded_assignments(assignments_list):
    """Combine selected graded assignments into one JSON for comparison (one per course).
    Each assignment can have content from: quiz questions, linked/attached files (PDF etc.), or written description.
    Multiple selected assignments are merged into a single combined_content string and one JSON."""
    if not assignments_list:
        return None
    parts = []
    index = []
    for idx, m in enumerate(assignments_list, 1):
        name = m.get('name', f'Assignment {idx}')
        content = m.get('content', '')
        sources = m.get('content_sources', [])
        index.append({
            'index': idx,
            'name': name,
            'id': m.get('id'),
            'points': m.get('points', 0),
            'content_sources': sources,
            'is_quiz': m.get('is_quiz', False),
            'had_linked_files': m.get('had_linked_files', False)
        })
        parts.append(f"\n{'='*80}\n📝 ASSIGNMENT {idx}: {name}\n{'='*80}\n{content}")
    combined_content = '\n\n'.join(parts)
    return {
        'id': 'graded_assignments_combined',
        'name': 'Selected Graded Materials (Combined)',
        'type': 'graded_assignments_combined',
        'total_assignments': len(index),
        'assignment_index': index,
        'combined_content': combined_content,
        'total_length': len(combined_content),
        'graded': True,
        'source_type': 'graded_assignments_combined'
    }


def combine_lecture_materials(files_list):
    """
    Combine selected lecture files into one JSON.
    For each file, if OCR produced slides use those; otherwise treat the
    whole extracted text as a single page.  Output format:

        [FileName - Slide 1] content
        [FileName - Slide 2] content
        ...
    """
    if not files_list:
        return None

    combined_parts = []
    file_index = []

    for idx, mat in enumerate(files_list, 1):
        name = mat.get('name', f'File {idx}')
        ocr = mat.get('ocr_result')

        slides_used = 0
        if ocr and ocr.get('slides'):
            for s in ocr['slides']:
                snum = s.get('slide_number', '')
                title = s.get('title', '')
                text = s.get('text', '')
                header = f"[{name} - Slide {snum}]"
                if title:
                    header += f" {title}"
                combined_parts.append(f"{header}\n{text}")
                slides_used += 1
        else:
            content = (mat.get('content') or '').strip()
            if content:
                combined_parts.append(f"[{name} - Page 1]\n{content}")
                slides_used = 1

        if slides_used:
            file_index.append({
                'index': idx,
                'name': name,
                'file_id': mat.get('id'),
                'file_type': mat.get('file_type', ''),
                'type': mat.get('source_type', 'unknown'),
                'slide_count': slides_used
            })

    if not combined_parts:
        return None

    combined_content = '\n\n'.join(combined_parts)

    result = {
        'id': 'lectures_combined',
        'name': 'Selected Lecture Materials (Combined)',
        'type': 'lectures_combined',
        'total_files': len(file_index),
        'file_index': file_index,
        'combined_content': combined_content,
        'total_length': len(combined_content),
        'graded': False,
        'source_type': 'combined_lectures'
    }

    print(f"✓ Combined {len(file_index)} lecture files into one JSON")
    print(f"  Total content length: {result['total_length']:,} characters")
    return result

def combine_discussions(discussions_list):
    """Combine selected discussion prompts into one JSON with section headers."""
    if not discussions_list:
        return None

    parts = []
    index = []
    for idx, d in enumerate(discussions_list, 1):
        name = d.get('name', f'Discussion {idx}')
        content = (d.get('content') or '').strip()
        author = d.get('author_name', '')
        posted_at = d.get('posted_at', '')
        if not content:
            continue
        index.append({
            'index': idx,
            'id': d.get('id'),
            'name': name,
            'author_name': author,
            'posted_at': posted_at
        })
        header = f"\n{'='*80}\n💬 DISCUSSION POST {idx}: {name}\n{'='*80}\n"
        meta_lines = []
        if author:
            meta_lines.append(f"Author: {author}")
        if posted_at:
            meta_lines.append(f"Posted At: {posted_at}")
        if meta_lines:
            header += "\n".join(meta_lines) + "\n\n"
        parts.append(f"{header}{content}")

    combined_content = '\n\n'.join(parts)
    return {
        'id': 'discussions_combined',
        'name': 'Selected Discussion Posts (Combined)',
        'type': 'discussions_combined',
        'total_discussions': len(index),
        'discussion_index': index,
        'combined_content': combined_content,
        'total_length': len(combined_content),
        'graded': False,
        'source_type': 'discussions_combined'
    }

def combine_all_selected_materials(files_list, assignments_list, discussions_list):
    """Combine selected lecture files, graded items, and discussions into one JSON."""
    parts = []
    index = []

    running_idx = 1
    for f in files_list or []:
        content = (f.get('content') or '').strip()
        if not content:
            continue
        name = f.get('name', f'File {running_idx}')
        index.append({'index': running_idx, 'type': 'lecture_file', 'id': f.get('id'), 'name': name})
        parts.append(f"\n{'='*80}\n📚 LECTURE FILE {running_idx}: {name}\n{'='*80}\n{content}")
        running_idx += 1

    for a in assignments_list or []:
        content = (a.get('content') or '').strip()
        if not content:
            continue
        name = a.get('name', f'Assignment {running_idx}')
        index.append({'index': running_idx, 'type': 'graded_material', 'id': a.get('id'), 'name': name})
        parts.append(f"\n{'='*80}\n📝 GRADED MATERIAL {running_idx}: {name}\n{'='*80}\n{content}")
        running_idx += 1

    for d in discussions_list or []:
        content = (d.get('content') or '').strip()
        if not content:
            continue
        name = d.get('name', f'Discussion {running_idx}')
        index.append({'index': running_idx, 'type': 'discussion_post', 'id': d.get('id'), 'name': name})
        parts.append(f"\n{'='*80}\n💬 DISCUSSION POST {running_idx}: {name}\n{'='*80}\n{content}")
        running_idx += 1

    if not parts:
        return None

    combined_content = '\n\n'.join(parts)
    return {
        'id': 'all_selected_combined',
        'name': 'All Selected Materials (Combined)',
        'type': 'all_selected_combined',
        'total_items': len(index),
        'item_index': index,
        'combined_content': combined_content,
        'total_length': len(combined_content),
        'source_type': 'all_selected_combined'
    }

def clean_html(html_content):
    """Convert HTML to clean plain text"""
    if not html_content:
        return ""
    
    try:
        soup = BeautifulSoup(html_content, 'html.parser')
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        # Get text
        text = soup.get_text()
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = ' '.join(chunk for chunk in chunks if chunk)
        return text
    except Exception as e:
        print(f"Error cleaning HTML: {e}")
        return html_content

def categorize_materials(materials, organized_materials):
    """Categorize materials into assignments, quizzes, labs, exams, files"""
    for material in materials:
        name = material['name'].lower()
        material_type = material['type']
        
        # Skip pages entirely
        if material_type == 'page':
            continue
        
        # Categorization logic
        if material_type == 'syllabus':
            # Syllabus gets its own category
            organized_materials['syllabus'].append(material)
        elif material_type == 'file':
            # All files go to lectures (slides, PDFs, docs)
            organized_materials['lectures'].append(material)
        elif material_type == 'quiz_question':
            # Individual quiz questions
            organized_materials['quiz_questions'].append(material)
        elif 'quiz' in name or material_type == 'quiz':
            organized_materials['quizzes'].append(material)
        elif 'lab' in name:
            organized_materials['labs'].append(material)
        elif 'exam' in name or 'test' in name or 'midterm' in name or 'final' in name:
            organized_materials['exams'].append(material)
        elif material_type == 'assignment':
            organized_materials['assignments'].append(material)
        elif material_type == 'discussion':
            organized_materials['discussions'].append(material)
        else:
            # Default: graded → assignments, ungraded → skip
            if material['graded']:
                organized_materials['assignments'].append(material)
            # Skip ungraded non-categorized materials

def save_course_data(course_id, data):
    """Save extracted course data to organized folder structure"""
    course_dir = os.path.join(DATA_DIR, f"course_{course_id}")
    
    os.makedirs(course_dir, exist_ok=True)
    material_type = data.get('material_type')

    # Save metadata
    metadata_path = os.path.join(course_dir, "metadata.json")
    with open(metadata_path, 'w', encoding='utf-8') as f:
        json.dump(data['metadata'], f, indent=2, ensure_ascii=False)
    
    # Save statistics
    stats_path = os.path.join(course_dir, "statistics.json")
    with open(stats_path, 'w', encoding='utf-8') as f:
        json.dump(data['statistics'], f, indent=2, ensure_ascii=False)
    
    # Syllabus: syllabus/ folder with syllabus_filtered.json and syllabus_raw.json
    if data['materials'].get('syllabus') and (material_type in ('syllabus', 'all') or not material_type):
        syllabus_dir = os.path.join(course_dir, "syllabus")
        os.makedirs(syllabus_dir, exist_ok=True)
        syllabus_data = data['materials']['syllabus'][0]
        slim_filtered = {
            'id': syllabus_data.get('id'),
            'name': syllabus_data.get('name'),
            'type': syllabus_data.get('type'),
            'filtered_data': syllabus_data.get('filtered_data'),
            'content': syllabus_data.get('content'),
            'graded': syllabus_data.get('graded', False),
            'points': syllabus_data.get('points', 0),
            'source_type': syllabus_data.get('source_type'),
            'source': syllabus_data.get('source'),
        }
        with open(os.path.join(syllabus_dir, "syllabus_filtered.json"), 'w', encoding='utf-8') as f:
            json.dump(slim_filtered, f, indent=2, ensure_ascii=False)
        if syllabus_data.get('raw_content'):
            raw_syllabus_data = {
                'id': f"syllabus_raw_{course_id}",
                'name': 'Course Syllabus (Raw)',
                'type': 'syllabus_raw',
                'content': syllabus_data['raw_content'],
                'graded': False,
                'points': 0,
                'source_type': 'syllabus_raw_untrimmed'
            }
            with open(os.path.join(syllabus_dir, "syllabus_raw.json"), 'w', encoding='utf-8') as f:
                json.dump(raw_syllabus_data, f, indent=2, ensure_ascii=False)
            print(f"  ✓ Saved syllabus/ (syllabus_filtered.json + syllabus_raw.json)")
        else:
            print(f"  ✓ Saved syllabus/syllabus_filtered.json")
    
    # Graded assignments (includes quizzes): graded_assignments/ folder with one combined JSON only
    # Include both assignments and quizzes lists—items with "quiz" in name were routed to quizzes by categorize_materials
    if material_type == 'graded_assignments':
        assignments_list = data['materials'].get('assignments', []) + data['materials'].get('quizzes', [])
        if assignments_list:
            graded_dir = os.path.join(course_dir, "graded_assignments")
            os.makedirs(graded_dir, exist_ok=True)
            combined = combine_graded_assignments(assignments_list)
            if combined:
                path = os.path.join(graded_dir, "graded_assignments_combined.json")
                with open(path, 'w', encoding='utf-8') as f:
                    json.dump(combined, f, indent=2, ensure_ascii=False)
                print(f"  ✓ Saved graded_assignments/graded_assignments_combined.json ({combined['total_assignments']} items)")
    elif material_type == 'all_selected':
        assignments_list = data['materials'].get('assignments', []) + data['materials'].get('quizzes', [])
        if assignments_list:
            graded_dir = os.path.join(course_dir, "graded_assignments")
            os.makedirs(graded_dir, exist_ok=True)
            combined = combine_graded_assignments(assignments_list)
            if combined:
                path = os.path.join(graded_dir, "graded_assignments_combined.json")
                with open(path, 'w', encoding='utf-8') as f:
                    json.dump(combined, f, indent=2, ensure_ascii=False)
                print(f"  ✓ Saved graded_assignments/graded_assignments_combined.json ({combined['total_assignments']} items)")

    # Lecture materials: lecture_materials/ folder with one combined JSON
    all_files = data['materials'].get('lectures', []) + data['materials'].get('labs', []) + data['materials'].get('exams', [])
    if all_files and material_type in ('lectures', 'all_selected'):
        lecture_dir = os.path.join(course_dir, "lecture_materials")
        os.makedirs(lecture_dir, exist_ok=True)
        combined_lectures = combine_lecture_materials(all_files)
        if combined_lectures:
            path = os.path.join(lecture_dir, "lectures_combined.json")
            with open(path, 'w', encoding='utf-8') as f:
                json.dump(combined_lectures, f, indent=2, ensure_ascii=False)
            print(f"  ✓ Saved lecture_materials/lectures_combined.json ({combined_lectures.get('total_files', 0)} files)")

    # Discussions: discussions/ folder with one combined JSON
    if material_type in ('discussions', 'all_selected'):
        discussions_list = data['materials'].get('discussions', [])
        if discussions_list:
            discussions_dir = os.path.join(course_dir, "discussions")
            os.makedirs(discussions_dir, exist_ok=True)
            combined_discussions = combine_discussions(discussions_list)
            if combined_discussions:
                path = os.path.join(discussions_dir, "discussions_combined.json")
                with open(path, 'w', encoding='utf-8') as f:
                    json.dump(combined_discussions, f, indent=2, ensure_ascii=False)
                print(f"  ✓ Saved discussions/discussions_combined.json ({combined_discussions.get('total_discussions', 0)} posts)")

    if material_type == 'all_selected':
        all_materials_dir = os.path.join(course_dir, "all_materials")
        os.makedirs(all_materials_dir, exist_ok=True)
        all_combined = combine_all_selected_materials(
            data['materials'].get('lectures', []),
            data['materials'].get('assignments', []) + data['materials'].get('quizzes', []),
            data['materials'].get('discussions', [])
        )
        if all_combined:
            path = os.path.join(all_materials_dir, "all_selected_combined.json")
            with open(path, 'w', encoding='utf-8') as f:
                json.dump(all_combined, f, indent=2, ensure_ascii=False)
            print(f"  ✓ Saved all_materials/all_selected_combined.json ({all_combined.get('total_items', 0)} items)")

def save_materials_by_type(course_dir, material_type, materials):
    """Save each material as a separate JSON file"""
    type_dir = os.path.join(course_dir, material_type)
    
    for idx, material in enumerate(materials, 1):
        # Syllabus is saved separately, so skip if encountered
        if material.get('type') == 'syllabus':
            continue
        
        filename = f"{material_type}_{idx:03d}.json"
        filepath = os.path.join(type_dir, filename)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(material, f, indent=2, ensure_ascii=False)

if __name__ == '__main__':
    # Log key routes so we can confirm /assignment-selection is registered
    rules = [r.rule for r in app.url_map.iter_rules() if 'assignment' in r.rule]
    if rules:
        print("Registered assignment routes:", rules)
    # Debug off by default for safer sharing; set FLASK_DEBUG=1 for local dev. No reloader: long extractions.
    _debug = os.environ.get("FLASK_DEBUG", "").lower() in ("1", "true", "yes")
    _port = int(os.environ.get("PORT", "5000"))
    app.run(debug=_debug, port=_port, use_reloader=False)
